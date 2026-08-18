#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
docflow_engine.py — 文档工作流引擎
===================================
用法:
  docflow_engine.py decode <out_path>          # stdin 为 base64 文本 → 解码写二进制文件
  docflow_engine.py extract <file>             # 提取文本 → stdout JSON {ok, format, text, chars, pages, slides, paragraphs, tables}
  docflow_engine.py create <fmt> <out_path>    # stdin 为 spec JSON → 生成精美文档 (docx/pptx/pdf/md/txt)
  docflow_engine.py edit   <fmt> <in> <out>    # stdin 为 spec JSON(ops) → 修改文档
  docflow_engine.py meta   <file>              # 元数据 → stdout JSON

spec JSON (create):
{
  "title": "…", "subtitle": "…", "author": "…", "date": "…",
  "theme": "blue|green|red|purple|gold|slate",
  "content": "# Markdown 文本…",              # 或
  "sections": [ {"type": "h1|h2|h3|p|bullets|numbered|table|quote|code|divider", …} ]
}

spec JSON (edit) — ops 数组:
  {"type":"replace","find":"…","replace":"…"}
  {"type":"set-title","title":"…"}
  {"type":"append","content":"…"} 或 {"type":"append","sections":[…]}
  {"type":"restyle","accent":"#RRGGBB"}
"""

import sys
import os
import io
import json
import base64
import re
import time
import datetime
import math


# 多用户路径栅栏根目录（None=管理员，不限路径；字符串=普通用户根，越界即拒绝）
FENCE_ROOT = None


def _wrap_lines(text, width_in, size):
    """按实际字符宽度估算文本在给定宽度（英寸）下折行后的行数。
    中文字符宽度 ≈ size/72 英寸，ASCII ≈ 0.52 倍；空文本按 1 行计。
    供 PPT 排版分页使用，避免固定高度导致的假性续页（续页后大面积空白）。"""
    try:
        width_in = float(width_in)
        size = float(size)
    except Exception:
        return 1
    if width_in <= 0:
        width_in = 1.0
    cjk = size / 72.0
    half = cjk * 0.52
    n = 0
    for seg in str(text or "").split("\n"):
        w = 0.0
        for ch in seg:
            w += cjk if ord(ch) > 0x2E7F else half
        n += max(int(math.ceil(w / width_in)), 1) if w > 0 else 1
    return max(n, 1)


def _fence(p):
    """普通用户模式下校验路径：必须位于 FENCE_ROOT 之内，否则拒绝。"""
    global FENCE_ROOT
    if FENCE_ROOT is None or p is None:
        return p
    try:
        rp = os.path.realpath(str(p))
    except Exception:
        rp = str(p)
    if rp == FENCE_ROOT or rp.startswith(FENCE_ROOT + os.sep):
        return p
    err("禁止访问用户目录之外的路径: %s" % p)

ROOT = os.path.dirname(os.path.abspath(__file__))

CJK_PDF_FONT = "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"
MONO_PDF_FONT = "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf"

# ---------------------------------------------------------------------------
# 主题配色
# ---------------------------------------------------------------------------
THEMES = {
    "blue":   {"accent": "#1F6FB2", "accent_dark": "#124F82", "light": "#E8F1FA", "soft": "#F5F9FD"},
    "green":  {"accent": "#2E8B57", "accent_dark": "#1E5F3C", "light": "#E6F4EC", "soft": "#F3FAF6"},
    "red":    {"accent": "#C0392B", "accent_dark": "#8E2A1F", "light": "#FBEAE8", "soft": "#FDF5F4"},
    "purple": {"accent": "#7D5BA6", "accent_dark": "#5A3F7E", "light": "#F0EAF7", "soft": "#F8F4FB"},
    "gold":   {"accent": "#B8860B", "accent_dark": "#85640A", "light": "#F9F0D8", "soft": "#FCF8EC"},
    "slate":  {"accent": "#455A64", "accent_dark": "#2C3A42", "light": "#E8EDEF", "soft": "#F4F7F8"},
}


def theme(name):
    return THEMES.get(name or "blue", THEMES["blue"])


# ---------------------------------------------------------------------------
# 基础工具
# ---------------------------------------------------------------------------
def err(msg):
    print(json.dumps({"ok": False, "error": str(msg)}, ensure_ascii=False))
    sys.exit(1)


def ok(data):
    print(json.dumps({"ok": True, **data}, ensure_ascii=False))


def ext_of(path):
    name = os.path.basename(path)
    m = re.search(r"\.([A-Za-z0-9]+)$", name)
    if not m:
        return ""
    e = m.group(1).lower()
    return {"markdown": "md"}.get(e, e)


def out(text, max_len=None):
    if max_len and len(text) > max_len:
        return text[:max_len] + "\n…[已截断]"
    return text


# ---------------------------------------------------------------------------
# 图片支持：下载 / 解析 / 搜索 / 识别
# ---------------------------------------------------------------------------
def _http_get_json(url, timeout=20):
    import urllib.request
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (compatible; docflow/1.0)"})
    with urllib.request.urlopen(req, timeout=timeout) as r:
        return json.loads(r.read().decode("utf-8", errors="replace"))


def resolve_image(src):
    """把图片 URL 或路径解析为本地可用的图片文件；网络图片会下载到 tmp。"""
    src = str(src or "").strip()
    if not src:
        return None
    if src.startswith(("http://", "https://")):
        import urllib.request
        import uuid
        ext = ".img"
        m = re.search(r"\.(png|jpe?g|gif|webp|svg)(?:\?|$)", src, re.I)
        if m:
            ext = "." + m.group(1).lower().replace("jpeg", "jpg")
        dest = os.path.join(ROOT, "tmp", "img_" + uuid.uuid4().hex + ext)
        try:
            req = urllib.request.Request(src, headers={"User-Agent": "Mozilla/5.0 (compatible; docflow/1.0)"})
            with urllib.request.urlopen(req, timeout=20) as r, open(dest, "wb") as f:
                f.write(r.read())
            return dest if os.path.exists(dest) else None
        except Exception:
            return None
    if os.path.exists(src):
        return src
    candidates = [
        os.path.join(os.path.dirname(ROOT), "uploads", src),
        os.path.join(os.path.dirname(ROOT), "outputs", src),
        os.path.join(ROOT, "tmp", src),
        os.path.join(os.path.dirname(ROOT), src),
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def image_search(query, max_results=5, image_type="all"):
    """通过 Wikimedia Commons 搜索图片/矢量图（无需 API Key）。"""
    import urllib.parse
    query = (query or "").strip()
    if not query:
        return {"error": "缺少搜索词"}
    n = min(max(int(max_results) if max_results else 5, 1), 20)
    params = {
        "action": "query",
        "format": "json",
        "generator": "search",
        "gsrsearch": query,
        "gsrnamespace": "6",
        "gsrlimit": str(n * 2),
        "prop": "imageinfo",
        "iiprop": "url|size|mime|extmetadata",
        "iiurlwidth": "800",
    }
    url = "https://commons.wikimedia.org/w/api.php?" + urllib.parse.urlencode(params)
    try:
        data = _http_get_json(url)
    except Exception as e:
        return {"error": "网络请求失败: %s" % e}
    pages = ((data.get("query") or {}).get("pages") or {})
    items = []
    for p in pages.values():
        title = p.get("title", "")
        info = (p.get("imageinfo") or [{}])[0]
        mime = info.get("mime", "")
        if image_type == "photo" and "svg" in mime:
            continue
        if image_type == "vector" and "svg" not in mime:
            continue
        items.append({
            "title": title,
            "url": info.get("url", ""),
            "thumb": info.get("thumburl", "") or info.get("url", ""),
            "width": info.get("width", 0),
            "height": info.get("height", 0),
            "mime": mime,
            "description": ((info.get("extmetadata") or {}).get("ImageDescription", {}) or {}).get("value", ""),
        })
        if len(items) >= n:
            break
    return {"items": items}


def image_info(path):
    from PIL import Image
    im = Image.open(path)
    info = {"format": im.format, "mode": im.mode, "width": im.width, "height": im.height}
    try:
        small = im.convert("RGB").resize((50, 50))
        pixels = list(small.getdata())
        n = max(len(pixels), 1)
        avg = tuple(sum(c[i] for c in pixels) // n for i in range(3))
        info["dominant_rgb"] = "#%02X%02X%02X" % avg
    except Exception:
        pass
    return info


def image_ocr(path):
    try:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(path), lang="chi_sim+eng").strip()
    except Exception:
        return ""


def image_recognize(path):
    """识别图片：先返回基础信息，再做 OCR；如配置视觉 API 则尝试多模态描述。"""
    if not os.path.exists(path):
        return {"error": "图片不存在: %s" % path}
    result = {"info": image_info(path), "ocr": image_ocr(path), "description": ""}
    # 可选：OpenAI 兼容视觉接口
    api_key = os.environ.get("DEEPSEEK_API_KEY") or os.environ.get("DOCFLOW_VISION_API_KEY")
    endpoint = os.environ.get("DOCFLOW_VISION_ENDPOINT") or os.environ.get("DEEPSEEK_BASE_URL")
    if api_key and endpoint:
        try:
            import base64
            import urllib.request
            with open(path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("ascii")
            mime = "image/png"
            low = path.lower()
            if low.endswith(".jpg") or low.endswith(".jpeg"):
                mime = "image/jpeg"
            elif low.endswith(".gif"):
                mime = "image/gif"
            elif low.endswith(".webp"):
                mime = "image/webp"
            payload = {
                "model": os.environ.get("DOCFLOW_VISION_MODEL", "deepseek-chat"),
                "messages": [
                    {"role": "user", "content": [
                        {"type": "text", "text": "请描述这张图片的内容，并给出可用于PPT配图的说明。"},
                        {"type": "image_url", "image_url": {"url": "data:%s;base64,%s" % (mime, b64)}},
                    ]}
                ],
                "max_tokens": 500,
            }
            req = urllib.request.Request(
                endpoint.rstrip("/") + "/chat/completions",
                data=json.dumps(payload).encode("utf-8"),
                headers={"Content-Type": "application/json", "Authorization": "Bearer " + api_key},
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=60) as r:
                data = json.loads(r.read().decode("utf-8"))
            result["description"] = ((data.get("choices") or [{}])[0].get("message") or {}).get("content", "")
        except Exception as e:
            result["vision_error"] = str(e)
    return result


# ---------------------------------------------------------------------------
# Markdown 解析 → sections
# ---------------------------------------------------------------------------
def parse_md(text):
    sections = []
    lines = (text or "").split("\n")
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i].rstrip()
        if not line.strip():
            i += 1
            continue
        if line.strip() in ("---", "***", "___"):
            sections.append({"type": "divider"})
            i += 1
            continue
        m = re.match(r"^(#{1,3})\s+(.*)$", line)
        if m:
            sections.append({"type": "h%d" % len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue
        m = re.match(r"^>\s?(.*)$", line)
        if m:
            quote = [m.group(1)]
            i += 1
            while i < n and lines[i].lstrip().startswith(">"):
                quote.append(lines[i].lstrip(">").strip())
                i += 1
            sections.append({"type": "quote", "text": " ".join(quote)})
            continue
        m = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", line)
        if m:
            sections.append({"type": "image", "src": m.group(2).strip(), "alt": m.group(1).strip()})
            i += 1
            continue
        if line.strip().startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            sections.append({"type": "code", "text": "\n".join(code)})
            continue
        if line.strip().startswith("|"):
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                    rows.append(cells)
                i += 1
            if rows:
                sections.append({"type": "table", "headers": rows[0], "rows": rows[1:]})
            continue
        m = re.match(r"^(\s*)[-*•]\s+(.*)$", line)
        if m:
            items = [m.group(2)]
            i += 1
            while i < n and re.match(r"^\s*[-*•]\s+", lines[i]):
                items.append(re.sub(r"^\s*[-*•]\s+", "", lines[i]).strip())
                i += 1
            sections.append({"type": "bullets", "items": items})
            continue
        m = re.match(r"^(\s*)\d+[\.、)]\s*(.*)$", line)
        if m:
            items = [m.group(2)]
            i += 1
            while i < n and re.match(r"^\s*\d+[\.、)]\s+", lines[i]):
                items.append(re.sub(r"^\s*\d+[\.、)]\s+", "", lines[i]).strip())
                i += 1
            sections.append({"type": "numbered", "items": items})
            continue
        para = [line.strip()]
        i += 1
        while i < n and lines[i].strip() and not re.match(
            r"^(#{1,3}\s|>\s?|```|\||\s*[-*•]\s+|\s*\d+[\.、)]\s+|---+$|\*\*\*+$)", lines[i]
        ):
            # 参考文献条目（[n] 开头）作为新段落起点，不并入上一段
            if re.match(r"^\[\d+\]", lines[i].strip()) and para:
                break
            para.append(lines[i].strip())
            i += 1
        sections.append({"type": "p", "text": " ".join(para)})
    return sections


def content_sections(spec):
    if spec.get("sections"):
        return spec["sections"]
    if spec.get("content"):
        return parse_md(spec["content"])
    return []


def spec_title(spec):
    return (spec.get("title") or "").strip()


# ===========================================================================
# DOCX
# ===========================================================================
_CITE_RE = __import__("re").compile(r"(\[\d+(?:[,\u2010\uFF0C\-]\d+)*\])")

def _cite_segments(text):
    """把正文中的 [n]、[n,m]、[n-m] 引用标注拆分为 (文本, 是否上标) 片段。"""
    segs = []
    pos = 0
    for m in _CITE_RE.finditer(text):
        if m.start() > pos:
            segs.append((text[pos:m.start()], False))
        # 段落开头的 [n]（后随空格+大写字母=文献列表条目）不上标
        sup = m.start() > 0
        segs.append((m.group(1), sup))
        pos = m.end()
    if pos < len(text):
        segs.append((text[pos:], False))
    return segs


def render_docx(spec, out_path):
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    t = theme(spec.get("theme"))
    doc = Document()
    for sec in doc.sections:
        sec.page_width = Cm(21.0)
        sec.page_height = Cm(29.7)
        sec.left_margin = Cm(2.6)
        sec.right_margin = Cm(2.6)
        sec.top_margin = Cm(2.6)
        sec.bottom_margin = Cm(2.4)

    normal = doc.styles["Normal"]
    normal.font.name = "微软雅黑"
    normal.font.size = Pt(11)
    try:
        rPr = normal.element.get_or_add_rPr()
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.append(rFonts)
        rFonts.set(qn("w:eastAsia"), "微软雅黑")
    except Exception:
        pass
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.4

    title = spec_title(spec)
    doc_title_text = title

    def add_par(text="", size=11, bold=False, italic=False, color=None, align=None,
                before=0, after=6, indent=None, hanging=None, line=1.4,
                border_bottom=None, font=None, east="微软雅黑"):
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(before)
        pf.space_after = Pt(after)
        pf.line_spacing = line
        if align is not None:
            p.alignment = align
        if indent is not None:
            pf.first_line_indent = Pt(indent)
        if hanging is not None:
            pf.left_indent = Pt(hanging)
            pf.first_line_indent = Pt(-hanging)
        if border_bottom:
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "12")
            bottom.set(qn("w:space"), "4")
            bottom.set(qn("w:color"), border_bottom.lstrip("#"))
            pBdr.append(bottom)
            pPr.append(pBdr)
        if text:
            # 引用数字角标上标：[n]、[n,m]、[n-m]（段落开头的 [n] 视为文献列表编号，不上标）
            for seg, sup in _cite_segments(text):
                r = p.add_run(seg)
                r.font.name = font or east
                if r._element.rPr is not None and r._element.rPr.rFonts is not None:
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), east)
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.italic = italic
                if sup:
                    r.font.superscript = True
                if color:
                    r.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        return p

    def shade_par(p, hexcolor):
        pPr = p._p.get_or_add_pPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hexcolor.lstrip("#"))
        pPr.append(shd)

    def shade_cell(cell, hexcolor):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hexcolor.lstrip("#"))
        tcPr.append(shd)

    def add_page_field(paragraph):
        run = paragraph.add_run()
        fld = OxmlElement("w:fldSimple")
        fld.set(qn("w:instr"), " PAGE ")
        rEl = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), "18")
        rPr.append(sz)
        col = OxmlElement("w:color")
        col.set(qn("w:val"), "888888")
        rPr.append(col)
        rEl.append(rPr)
        tEl = OxmlElement("w:t")
        tEl.text = "1"
        rEl.append(tEl)
        fld.append(rEl)
        paragraph._p.append(fld)

    # ---- 页眉/页脚 ----
    if title:
        try:
            hp = doc.sections[0].header.paragraphs[0]
            hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            hr = hp.add_run(title)
            hr.font.size = Pt(9)
            hr.font.color.rgb = RGBColor.from_string("999999")
            hr.font.name = "微软雅黑"
            if hr._element.rPr is not None and hr._element.rPr.rFonts is not None:
                hr._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
        except Exception:
            pass
    try:
        fp = doc.sections[0].footer.paragraphs[0]
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1 = fp.add_run("第 ")
        r1.font.size = Pt(9)
        r1.font.color.rgb = RGBColor.from_string("888888")
        add_page_field(fp)
        r2 = fp.add_run(" 页")
        r2.font.size = Pt(9)
        r2.font.color.rgb = RGBColor.from_string("888888")
    except Exception:
        pass

    # ---- 封面 ----
    if title:
        add_par(title, size=26, bold=True, color=t["accent_dark"],
                align=WD_ALIGN_PARAGRAPH.CENTER, before=0, after=12, line=1.25)
        add_par("", after=10, border_bottom=t["accent"])
        if spec.get("subtitle"):
            add_par(spec["subtitle"], size=14, color="#666666",
                    align=WD_ALIGN_PARAGRAPH.CENTER, after=6)
        meta = []
        if spec.get("author"):
            meta.append(spec["author"])
        if spec.get("date"):
            meta.append(spec["date"])
        if meta:
            add_par("　".join(meta), size=11, color="#888888",
                    align=WD_ALIGN_PARAGRAPH.CENTER, after=18)
        add_par("", after=8)

    # ---- 正文渲染 ----
    def render_section(s):
        typ = s.get("type")
        if typ == "h1":
            add_par(s.get("text", ""), size=17, bold=True, color=t["accent"],
                    before=18, after=8, border_bottom=t["accent"], line=1.25)
        elif typ == "h2":
            add_par(s.get("text", ""), size=14, bold=True, color=t["accent_dark"],
                    before=14, after=6, line=1.25)
        elif typ == "h3":
            add_par(s.get("text", ""), size=12, bold=True, color="#444444",
                    before=10, after=4, line=1.25)
        elif typ == "p":
            add_par(s.get("text", ""), indent=22)
        elif typ == "quote":
            add_par(s.get("text", ""), italic=True, color="#666666",
                    before=4, after=10, indent=24, line=1.35)
        elif typ == "code":
            for line in (s.get("text", "") or "").split("\n"):
                p = add_par(line if line else " ", size=9.5, after=0, line=1.2,
                            font="Consolas", east="宋体")
                pf = p.paragraph_format
                pf.left_indent = Pt(16)
                pf.right_indent = Pt(16)
                shade_par(p, t["light"])
            add_par("", size=4, after=8)
        elif typ == "bullets":
            for it in s.get("items", []):
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = Pt(20)
                pf.first_line_indent = Pt(-14)
                pf.space_after = Pt(4)
                pf.line_spacing = 1.35
                rb = p.add_run("•  ")
                rb.font.size = Pt(11)
                rb.font.bold = True
                rb.font.color.rgb = RGBColor.from_string(t["accent"].lstrip("#"))
                rb.font.name = "微软雅黑"
                if rb._element.rPr is not None and rb._element.rPr.rFonts is not None:
                    rb._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                for _seg, _sup in _cite_segments(str(it)):
                    rt = p.add_run(_seg)
                    rt.font.size = Pt(11)
                    rt.font.name = "微软雅黑"
                    if rt._element.rPr is not None and rt._element.rPr.rFonts is not None:
                        rt._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                    if _sup:
                        rt.font.superscript = True
        elif typ == "numbered":
            for n, it in enumerate(s.get("items", []), 1):
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = Pt(22)
                pf.first_line_indent = Pt(-16)
                pf.space_after = Pt(4)
                pf.line_spacing = 1.35
                rn = p.add_run("%d. " % n)
                rn.font.size = Pt(11)
                rn.font.bold = True
                rn.font.color.rgb = RGBColor.from_string(t["accent"].lstrip("#"))
                rn.font.name = "微软雅黑"
                if rn._element.rPr is not None and rn._element.rPr.rFonts is not None:
                    rn._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                for _seg, _sup in _cite_segments(str(it)):
                    rt = p.add_run(_seg)
                    rt.font.size = Pt(11)
                    rt.font.name = "微软雅黑"
                    if rt._element.rPr is not None and rt._element.rPr.rFonts is not None:
                        rt._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                    if _sup:
                        rt.font.superscript = True
        elif typ == "table":
            render_table(s)
        elif typ == "divider":
            add_par("", after=8, border_bottom="#CCCCCC")
        elif typ == "image":
            local = resolve_image(s.get("src", ""))
            if not local:
                add_par("图片不可用: " + str(s.get("src", "")), italic=True, color="#999999", after=8)
                return
            try:
                from PIL import Image as _Img
                _im = _Img.open(local)
                _iw, _ih = _im.size
                # 页面可用宽度约 16cm，图片最大宽 15cm、最大高 12cm，等比缩放
                _w = 15.0
                if _iw and _ih:
                    _h = _w * (_ih / float(_iw))
                    if _h > 12.0:
                        _h = 12.0
                        _w = _h * (_iw / float(_ih))
                else:
                    _h = None
                _p = doc.add_paragraph()
                _p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                _p.paragraph_format.space_before = Pt(6)
                _p.paragraph_format.space_after = Pt(2)
                _r = _p.add_run()
                if _h:
                    _r.add_picture(local, width=Cm(_w), height=Cm(_h))
                else:
                    _r.add_picture(local, width=Cm(_w))
                _alt = s.get("alt", "")
                if _alt:
                    add_par(_alt, size=9, color="#888888", align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
            except Exception as e:
                add_par("图片加载失败: %s" % e, italic=True, color="#999999", after=8)

    def render_table(s):
        headers = s.get("headers") or []
        rows = s.get("rows") or []
        ncols = max(len(headers), max((len(r) for r in rows), default=0), 1)
        table = doc.add_table(rows=1 + len(rows), cols=ncols)
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for j in range(ncols):
            cell = table.rows[0].cells[j]
            cell.text = ""
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(headers[j] if j < len(headers) else "")
            r.font.bold = True
            r.font.size = Pt(10.5)
            r.font.color.rgb = RGBColor.from_string("FFFFFF")
            r.font.name = "微软雅黑"
            if r._element.rPr is not None and r._element.rPr.rFonts is not None:
                r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
            shade_cell(cell, t["accent"])
        for i, row in enumerate(rows):
            for j in range(ncols):
                cell = table.rows[i + 1].cells[j]
                cell.text = ""
                p = cell.paragraphs[0]
                r = p.add_run(str(row[j]) if j < len(row) else "")
                r.font.size = Pt(10)
                r.font.name = "微软雅黑"
                if r._element.rPr is not None and r._element.rPr.rFonts is not None:
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                if i % 2 == 1:
                    shade_cell(cell, t["soft"])
        add_par("", size=4, after=10)

    for s in content_sections(spec):
        render_section(s)

    # 主题元信息（供 restyle 识别）
    try:
        core = doc.core_properties
        core.title = title or "文档"
        core.author = spec.get("author") or "docflow"
    except Exception:
        pass

    doc.save(out_path)


def replace_in_paragraph(p, find, repl):
    full = "".join(r.text for r in p.runs)
    if find not in full:
        return False
    new = full.replace(find, repl)
    if p.runs:
        p.runs[0].text = new
        for r in p.runs[1:]:
            r.text = ""
    return True


def edit_docx(spec, in_path, out_path):
    from docx import Document
    from docx.shared import RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    t = theme(spec.get("theme"))
    doc = Document(in_path)
    ops = spec.get("ops") or []

    def cell_shade_fill(cell):
        tcPr = cell._tc.find(qn("w:tcPr"))
        if tcPr is None:
            return None
        shd = tcPr.find(qn("w:shd"))
        if shd is None:
            return None
        return shd.get(qn("w:fill"))

    for op in ops:
        typ = op.get("type")
        if typ == "replace":
            find = op.get("find", "")
            repl = op.get("replace", "")
            if not find:
                continue
            for p in doc.paragraphs:
                replace_in_paragraph(p, find, repl)
            for tb in doc.tables:
                for row in tb.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            replace_in_paragraph(p, find, repl)
        elif typ == "set-title":
            title = (op.get("title") or "").strip()
            if not title:
                continue
            target = None
            for p in doc.paragraphs:
                if p.style is not None and "heading" in p.style.name.lower():
                    target = p
                    break
            if target is None and doc.paragraphs:
                target = doc.paragraphs[0]
            if target is not None:
                if target.runs:
                    target.runs[0].text = title
                    for r in target.runs[1:]:
                        r.text = ""
                else:
                    target.add_run(title)
        elif typ == "append":
            for s in content_sections(op):
                _docx_append_section(doc, s, t)
        elif typ == "restyle":
            accent = (op.get("accent") or "").strip()
            if not re.fullmatch(r"#[0-9A-Fa-f]{6}", accent):
                continue
            new_t = theme("blue")
            new_t["accent"] = accent
            new_t["accent_dark"] = accent
            old_colors = set()
            for th in THEMES.values():
                old_colors.add(th["accent"].lstrip("#").upper())
                old_colors.add(th["accent_dark"].lstrip("#").upper())
            for p in doc.paragraphs:
                for r in p.runs:
                    try:
                        c = r.font.color
                        if c is not None and c.rgb is not None:
                            s = str(c.rgb).upper()
                            if s in old_colors:
                                c.rgb = RGBColor.from_string(accent.lstrip("#"))
                    except Exception:
                        pass
            for tb in doc.tables:
                for row in tb.rows:
                    for cell in row.cells:
                        f = cell_shade_fill(cell)
                        if f is not None and f.upper() in old_colors:
                            tcPr = cell._tc.get_or_add_tcPr()
                            shd = tcPr.find(qn("w:shd"))
                            if shd is not None:
                                shd.set(qn("w:fill"), accent.lstrip("#"))
                        for p in cell.paragraphs:
                            for r in p.runs:
                                try:
                                    c = r.font.color
                                    if c is not None and c.rgb is not None:
                                        s = str(c.rgb).upper()
                                        if s in old_colors:
                                            c.rgb = RGBColor.from_string(accent.lstrip("#"))
                                except Exception:
                                    pass
    doc.save(out_path)


def _docx_append_section(doc, s, t):
    """在已有 docx 末尾追加一个 section（复用渲染逻辑的简化版）。"""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    def add_par(text="", size=11, bold=False, italic=False, color=None, align=None,
                before=0, after=6, indent=None, hanging=None, line=1.4, border_bottom=None):
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(before)
        pf.space_after = Pt(after)
        pf.line_spacing = line
        if align is not None:
            p.alignment = align
        if indent is not None:
            pf.first_line_indent = Pt(indent)
        if hanging is not None:
            pf.left_indent = Pt(hanging)
            pf.first_line_indent = Pt(-hanging)
        if border_bottom:
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "12")
            bottom.set(qn("w:space"), "4")
            bottom.set(qn("w:color"), border_bottom.lstrip("#"))
            pBdr.append(bottom)
            pPr.append(pBdr)
        if text:
            for seg, sup in _cite_segments(text):
                r = p.add_run(seg)
                r.font.name = "微软雅黑"
                if r._element.rPr is not None and r._element.rPr.rFonts is not None:
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.italic = italic
                if sup:
                    r.font.superscript = True
                if color:
                    r.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        return p

    def shade_cell(cell, hexcolor):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hexcolor.lstrip("#"))
        tcPr.append(shd)

    typ = s.get("type")
    if typ == "h1":
        add_par(s.get("text", ""), size=17, bold=True, color=t["accent"], before=18, after=8, border_bottom=t["accent"])
    elif typ == "h2":
        add_par(s.get("text", ""), size=14, bold=True, color=t["accent_dark"], before=14, after=6)
    elif typ == "h3":
        add_par(s.get("text", ""), size=12, bold=True, color="#444444", before=10, after=4)
    elif typ == "image":
        from docx.shared import Cm as _Cm
        local = resolve_image(s.get("src", ""))
        if not local:
            add_par("图片不可用: " + str(s.get("src", "")), italic=True, color="#999999", after=8)
            return
        try:
            from PIL import Image as _Img
            _im = _Img.open(local)
            _iw, _ih = _im.size
            _w = 15.0
            if _iw and _ih:
                _h = _w * (_ih / float(_iw))
                if _h > 12.0:
                    _h = 12.0
                    _w = _h * (_iw / float(_ih))
            else:
                _h = None
            _p = doc.add_paragraph()
            _p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _p.paragraph_format.space_before = Pt(6)
            _p.paragraph_format.space_after = Pt(2)
            _r = _p.add_run()
            if _h:
                _r.add_picture(local, width=_Cm(_w), height=_Cm(_h))
            else:
                _r.add_picture(local, width=_Cm(_w))
            _alt = s.get("alt", "")
            if _alt:
                add_par(_alt, size=9, color="#888888", align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        except Exception as e:
            add_par("图片加载失败: %s" % e, italic=True, color="#999999", after=8)
    elif typ == "p":
        add_par(s.get("text", ""), indent=22)
    elif typ == "quote":
        add_par(s.get("text", ""), italic=True, color="#666666", before=4, after=10, indent=24)
    elif typ == "code":
        for line in (s.get("text", "") or "").split("\n"):
            p = add_par(line if line else " ", size=9.5, after=0, line=1.2)
            pf = p.paragraph_format
            pf.left_indent = Pt(16)
            pf.right_indent = Pt(16)
            pPr = p._p.get_or_add_pPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), t["light"].lstrip("#"))
            pPr.append(shd)
        add_par("", size=4, after=8)
    elif typ == "bullets":
        for it in s.get("items", []):
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.left_indent = Pt(20)
            pf.first_line_indent = Pt(-14)
            pf.space_after = Pt(4)
            pf.line_spacing = 1.35
            rb = p.add_run("•  ")
            rb.font.size = Pt(11)
            rb.font.bold = True
            rb.font.color.rgb = RGBColor.from_string(t["accent"].lstrip("#"))
            rb.font.name = "微软雅黑"
            for _seg, _sup in _cite_segments(str(it)):
                rt = p.add_run(_seg)
                rt.font.size = Pt(11)
                rt.font.name = "微软雅黑"
                if _sup:
                    rt.font.superscript = True
    elif typ == "numbered":
        for n, it in enumerate(s.get("items", []), 1):
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.left_indent = Pt(22)
            pf.first_line_indent = Pt(-16)
            pf.space_after = Pt(4)
            pf.line_spacing = 1.35
            rn = p.add_run("%d. " % n)
            rn.font.size = Pt(11)
            rn.font.bold = True
            rn.font.color.rgb = RGBColor.from_string(t["accent"].lstrip("#"))
            rn.font.name = "微软雅黑"
            for _seg, _sup in _cite_segments(str(it)):
                rt = p.add_run(_seg)
                rt.font.size = Pt(11)
                rt.font.name = "微软雅黑"
                if _sup:
                    rt.font.superscript = True
    elif typ == "table":
        headers = s.get("headers") or []
        rows = s.get("rows") or []
        ncols = max(len(headers), max((len(r) for r in rows), default=0), 1)
        table = doc.add_table(rows=1 + len(rows), cols=ncols)
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for j in range(ncols):
            cell = table.rows[0].cells[j]
            cell.text = ""
            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(headers[j] if j < len(headers) else "")
            r.font.bold = True
            r.font.size = Pt(10.5)
            r.font.color.rgb = RGBColor.from_string("FFFFFF")
            r.font.name = "微软雅黑"
            shade_cell(cell, t["accent"])
        for i, row in enumerate(rows):
            for j in range(ncols):
                cell = table.rows[i + 1].cells[j]
                cell.text = ""
                p = cell.paragraphs[0]
                r = p.add_run(str(row[j]) if j < len(row) else "")
                r.font.size = Pt(10)
                r.font.name = "微软雅黑"
                if i % 2 == 1:
                    shade_cell(cell, t["soft"])
        add_par("", size=4, after=10)
    elif typ == "divider":
        add_par("", after=8, border_bottom="#CCCCCC")


def extract_docx(path):
    from docx.oxml.ns import qn
    from docx import Document

    doc = Document(path)
    parts = []
    tables = 0
    paras = 0
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            paras += 1
            pstyle = child.find(qn("w:pPr"))
            style_name = ""
            if pstyle is not None:
                ps = pstyle.find(qn("w:pStyle"))
                if ps is not None:
                    style_name = ps.get(qn("w:val")) or ""
            text = "".join(t.text or "" for t in child.iter(qn("w:t")))
            m = re.match(r"^(Heading|标题)\s*(\d)?", style_name)
            if m:
                lvl = m.group(2) or "1"
                parts.append("#" * int(lvl) + " " + text.strip())
            elif text.strip():
                parts.append(text.strip())
        elif child.tag == qn("w:tbl"):
            tables += 1
            for row in child.iter(qn("w:tr")):
                cells = []
                for tc in row.iter(qn("w:tc")):
                    t = "".join(x.text or "" for x in tc.iter(qn("w:t")))
                    cells.append(t.strip())
                parts.append(" | ".join(cells))
    return {
        "format": "docx",
        "text": "\n".join(p for p in parts if p),
        "paragraphs": paras,
        "tables": tables,
    }


def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    paras = 0
    tables = 0

    def walk(shapes):
        nonlocal paras, tables
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                walk(shape.shapes)
                continue
            if shape.has_table:
                tables += 1
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
            elif shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs).strip()
                    if txt:
                        paras += 1
                        parts.append(txt)

    slide_count = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_count = i
        parts.append("【第 %d 页】" % i)
        walk(slide.shapes)
    return {"format": "pptx", "text": "\n".join(parts), "slides": slide_count, "paragraphs": paras, "tables": tables}


def extract_pdf(path):
    import pdfplumber

    parts = []
    pages = 0
    with pdfplumber.open(path) as pdf:
        pages = len(pdf.pages)
        for page in pdf.pages:
            txt = page.extract_text() or ""
            parts.append(txt)
    return {"format": "pdf", "text": "\n".join(parts), "pages": pages}


def extract_text_file(path):
    with io.open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    lines = text.split("\n")
    return {"format": ext_of(path), "text": text, "paragraphs": len([l for l in lines if l.strip()])}


def extract_xlsx(path):
    """提取 Excel：每个工作表输出为表格行（| 分隔），含 sheet 标题。"""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    sheets = 0
    rows = 0
    for ws in wb.worksheets:
        sheets += 1
        parts.append("【工作表：%s】" % (ws.title or "Sheet%d" % sheets))
        for row in ws.iter_rows(values_only=True):
            vals = []
            has = False
            for c in row:
                if c is None:
                    vals.append("")
                else:
                    has = True
                    s = str(c).strip()
                    vals.append(s)
            if has and any(vals):
                rows += 1
                parts.append(" | ".join(vals).rstrip(" |"))
    wb.close()
    return {"format": "xlsx", "text": "\n".join(parts), "sheets": sheets, "rows": rows}


def extract_image(path):
    try:
        info = image_info(path)
        ocr = image_ocr(path)
        text = "图片信息：格式 %s，尺寸 %dx%d，主色 %s" % (
            info.get("format", ""), info.get("width", 0), info.get("height", 0), info.get("dominant_rgb", ""))
        if ocr:
            text += "\nOCR识别：\n" + ocr
        return {"format": ext_of(path) or "image", "text": text, "image": info}
    except Exception as e:
        return {"format": ext_of(path) or "image", "text": "", "error": str(e)}


def extract_any(path):
    fmt = ext_of(path)
    if fmt == "docx":
        return extract_docx(path)
    if fmt == "pptx":
        return extract_pptx(path)
    if fmt == "pdf":
        return extract_pdf(path)
    if fmt == "xlsx":
        return extract_xlsx(path)
    if fmt in ("txt", "md"):
        return extract_text_file(path)
    if fmt in ("png", "jpg", "jpeg", "gif", "webp", "bmp", "svg"):
        return extract_image(path)
    return {"format": fmt, "text": ""}


# ===========================================================================
# PPTX
# ===========================================================================
def render_pptx_visual(spec, out_path):
    """可视化大字版 PPT 渲染器：正文最小 24pt，卡片网格铺满页面，
    大色块/编号圆点/图标字符生动呈现。风格借鉴「个人智能体部署指南」示例。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree

    t = theme(spec.get("theme"))

    def C(h):
        return RGBColor.from_string(str(h).lstrip("#"))

    accent = C(t["accent"])
    accent_dark = C(t["accent_dark"])
    light = C(t["light"])
    soft = C(t["soft"])
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    MED_GRAY = RGBColor(0x66, 0x66, 0x66)
    LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    BG_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

    def hx(c):
        return "%02X%02X%02X" % (c[0], c[1], c[2])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = 13.333, 7.5
    LM = 0.8
    CW = 11.733
    slides = []

    def style_run(r, size=24, bold=False, color="333333", font="微软雅黑", italic=False):
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = RGBColor.from_string(str(color).lstrip("#"))
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)

    def clean(shape):
        sp = shape._element
        st = sp.find(qn("p:style"))
        if st is not None:
            sp.remove(st)

    def add_rect(slide, l, tp, w, h, color, radius=None):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = slide.shapes.add_shape(shape_type, Inches(l), Inches(tp), Inches(w), Inches(h))
        if radius:
            try:
                sp.adjustments[0] = radius
            except Exception:
                pass
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        clean(sp)
        return sp

    def add_hline(slide, x, y, length, color, thickness=0.5):
        h = max(thickness * 12700, 6350)
        return add_rect(slide, x, y, length, h / 914400.0, color)

    def add_oval(slide, l, tp, size, letter, bg, fg=WHITE, font_size=24):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(tp), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = bg
        c.line.fill.background()
        c.shadow.inherit = False
        clean(c)
        tf = c.text_frame
        tf.paragraphs[0].text = letter
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.set("anchor", "ctr")
        for a in ["lIns", "tIns", "rIns", "bIns"]:
            bodyPr.set(a, "0")
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = fg
            r.font.name = "微软雅黑"
        return c

    def add_box(slide, l, tp, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    def add_text(slide, l, tp, w, h, text, size=24, bold=False, color=DARK_GRAY,
                 font="微软雅黑", align=PP_ALIGN.LEFT, anchor="t"):
        tb = slide.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.set("anchor", anchor)
        for a in ["lIns", "tIns", "rIns", "bIns"]:
            bodyPr.set(a, "45720")
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if i > 0:
                p.space_before = Pt(4)
            r = p.add_run()
            r.text = line
            style_run(r, size, bold, hx(color), font)
        return tb

    def add_slide():
        s = prs.slides.add_slide(blank)
        slides.append(s)
        # 每页右上角加一个小图标，保证“每页一图”
        icon_path = os.path.join(os.path.dirname(ROOT), 'assets', 'dsh_icon.png')
        local_icon = resolve_image(icon_path)
        if local_icon:
            try:
                from PIL import Image
                im = Image.open(local_icon)
                iw, ih = im.size
                if iw and ih:
                    icon_h = 0.55
                    icon_w = icon_h * iw / ih
                    s.shapes.add_picture(local_icon, Inches(W - icon_w - 0.25), Inches(0.18),
                                         width=Inches(icon_w), height=Inches(icon_h))
            except Exception:
                pass
        return s

    def est_lines(text, width_in, size=24):
        """估算在给定宽度（英寸）下的折行行数（中文字≈size/72 英寸宽，ASCII≈0.52倍）。"""
        return _wrap_lines(text, width_in, size)

    # ---------------- 封面：左侧大色块 + 右侧大标题 ----------------
    s = add_slide()
    add_rect(s, 0, 0, 5.0, H, accent)
    add_rect(s, 5.0, 0, 0.14, H, accent_dark)
    add_oval(s, 3.6, 5.6, 2.6, "", light, fg=light)   # 装饰圆
    add_oval(s, 4.2, 6.3, 1.1, "", WHITE, fg=WHITE)   # 装饰圆
    brand = (spec.get("author") or "").split("·")[0].strip()
    if brand:
        add_text(s, 0.55, 0.5, 3.9, 0.5, brand, size=24, bold=True, color=WHITE)
    title = spec_title(spec) or "文档标题"
    lines = title.split("\n") if isinstance(title, str) else [title]
    ty = 2.0
    for ln in lines:
        add_text(s, 0.55, ty, 4.2, 0.95, ln, size=40, bold=True, color=WHITE)
        ty += 0.95
    if spec.get("subtitle"):
        add_text(s, 5.6, 2.1, 7.3, 1.1, spec["subtitle"], size=28, bold=True, color=accent_dark)
    add_hline(s, 5.6, 3.35, 2.2, accent, 2.5)
    body_brief = (spec.get("content") or "")[:120]
    add_text(s, 5.6, 3.7, 7.0, 1.2, "面向新手 · 可视化讲解 · 动手即会", size=24, color=MED_GRAY)
    meta = []
    if spec.get("author"):
        meta.append(spec["author"])
    if spec.get("date"):
        meta.append(spec["date"])
    if meta:
        add_text(s, 5.6, 6.35, 7.0, 0.6, "　".join(meta), size=24, color=MED_GRAY)

    # ---------------- 内容结构 ----------------
    cur = None
    cur_title = ""
    body_tf = None
    est_h = 0.0
    BODY_TOP = 1.95
    BODY_MAX = 5.0
    h2_num = 0
    part_num = 0

    def open_slide(title_text, is_divider=False, nav=None):
        nonlocal cur, cur_title, body_tf, est_h, part_num
        cur = add_slide()
        cur_title = title_text
        est_h = 0.0
        if is_divider:
            part_num += 1
            add_rect(cur, 0, 0, 0.55, H, accent)
            add_oval(cur, 0.9, 1.5, 1.7, "", light, fg=light)
            add_text(cur, 1.05, 1.75, 3.0, 1.2, "PART %02d" % part_num, size=44, bold=True, color=accent)
            add_text(cur, 1.05, 3.1, 11.2, 1.6, title_text, size=44, bold=True, color=accent_dark)
            add_hline(cur, 1.1, 5.0, 3.5, accent, 2.5)
            # 章节导航：列出本章小节，避免过渡页大块空白
            if nav:
                add_text(cur, 1.1, 5.35, 11.0, 0.5, "本章内容", size=20, bold=True, color=MED_GRAY)
                ny = 5.95
                for i, h2t in enumerate(nav[:6]):
                    col = i % 2
                    row = i // 2
                    add_text(cur, 1.1 + col * 5.8, ny + row * 0.55, 5.6, 0.5,
                             "· " + h2t, size=20, color=accent_dark)
            body_tf = None
        else:
            add_text(cur, LM, 0.5, 11.7, 0.9, title_text, size=32, bold=True, color=accent_dark, anchor="b")
            add_hline(cur, LM, 1.52, 1.9, accent, 3.0)
            add_hline(cur, LM, 1.52, 11.7, LINE_GRAY, 0.5)
            body_tf = add_box(cur, LM, BODY_TOP, CW, BODY_MAX)
        return cur

    def ensure_body():
        nonlocal cur
        if body_tf is None:
            open_slide(cur_title)

    cont_n = {}

    def ensure_space(need):
        nonlocal est_h
        # 只有确实放不下才续页（+0.1in 容差），避免短段落触发假性续页
        if est_h + need > BODY_MAX + 0.1:
            n = cont_n.get(cur_title, 0) + 1
            cont_n[cur_title] = n
            open_slide(cur_title + ("（续%d）" % n if n > 1 else "（续）"))

    def para(text="", size=24, bold=False, color="333333", before=0, after=8):
        nonlocal est_h
        ensure_body()
        n_lines = est_lines(text, CW, size)
        lh = 0.28 if size <= 16 else 0.46
        need = n_lines * lh + 0.12 + after / 72.0
        ensure_space(need)
        p = body_tf.paragraphs[0] if not body_tf.paragraphs[0].runs and est_h == 0 else body_tf.add_paragraph()
        p.space_before = Pt(before)
        p.space_after = Pt(after)
        r = p.add_run()
        r.text = text
        style_run(r, size, bold, color)
        est_h += need

    def add_h2_row(text):
        nonlocal est_h, h2_num
        ensure_body()
        ensure_space(0.75)
        h2_num += 1
        y = BODY_TOP + est_h
        add_oval(cur, LM, y - 0.05, 0.6, str(h2_num), accent, font_size=26)
        tf = add_box(cur, LM + 0.85, y - 0.1, CW - 0.85, 0.62)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style_run(r, 28, True, hx(accent_dark))
        est_h += 0.75

    def add_bullet_cards(items):
        """要点卡片网格：2 列圆角卡片，24pt 文字，按高度逐行分批铺满内容区。
        修复：整批 4 卡超高时不再整体续页掏空标题页，而是按行拆分，
        每页至少放下一行卡片；末行单卡时使用整行宽大卡片填满视觉。"""
        nonlocal est_h
        ensure_body()
        n = len(items)
        idx = 0
        gap = 0.35
        cw2 = (CW - gap) / 2.0          # 2 列时单卡宽度

        def _card_h(it):
            nl = est_lines(it, cw2 - 0.7, 24)
            return min(max(nl * 0.46 + 0.72, 1.3), 2.55)

        while idx < n:
            # 逐行收集（每行最多 2 卡，行高取 max），最多收集 2 行
            rows_info = []              # [(row_h, [item_idx,...]), ...]
            i = idx
            while i < n and len(rows_info) < 2:
                h1 = _card_h(items[i])
                j = i + 1
                if j < n:
                    h2 = _card_h(items[j])
                    rows_info.append((max(h1, h2), [i, j]))
                    i = j + 1
                else:
                    rows_info.append((h1, [i]))
                    i = j
            # 在剩余空间内确定可放行数（至少 1 行）
            take = len(rows_info)
            while take > 1:
                tot = sum(r[0] for r in rows_info[:take]) + gap * (take - 1)
                if est_h + tot <= BODY_MAX + 0.02:
                    break
                take -= 1
            batch_rows = rows_info[:take]
            total = sum(r[0] for r in batch_rows) + gap * (take - 1)
            ensure_space(total + 0.1)
            y = BODY_TOP + est_h
            for ri, (rh, ids) in enumerate(batch_rows):
                single = len(ids) == 1
                cw = CW if single else cw2
                for k, it_idx in enumerate(ids):
                    x = LM if single else LM + k * (cw + gap)
                    add_rect(cur, x, y, cw, rh, light, radius=0.12)
                    add_rect(cur, x, y, 0.12, rh, accent)
                    add_text(cur, x + 0.35, y + 0.12, cw - 0.7, 0.5, "▪", size=30, bold=True, color=accent)
                    tf = add_box(cur, x + 0.35, y + 0.64, cw - 0.7, rh - 0.74)
                    p = tf.paragraphs[0]
                    r_2 = p.add_run()
                    r_2.text = items[it_idx]
                    style_run(r_2, 24, False, hx(DARK_GRAY))
                y += rh + gap
            est_h += total + gap
            idx += sum(len(r[1]) for r in batch_rows)

    def add_numbered(items):
        nonlocal est_h
        ensure_body()
        for i, it in enumerate(items):
            num = i + 1
            nl = est_lines(it, CW - 1.2, 24)
            need = max(nl * 0.44 + 0.05, 0.66)
            ensure_space(need)
            y = BODY_TOP + est_h
            add_oval(cur, LM, y, 0.55, str(num), accent, font_size=24)
            tf = add_box(cur, LM + 0.85, y - 0.05, CW - 0.85, need)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = it
            style_run(r, 24, False, hx(DARK_GRAY))
            est_h += need + 0.14

    def add_quote(text):
        nonlocal est_h
        ensure_body()
        nl = est_lines(text, CW - 1.4, 26)
        hh = max(nl * 0.46 + 0.45, 0.9)
        ensure_space(hh + 0.1)
        y = BODY_TOP + est_h
        add_rect(cur, LM, y, CW, hh, accent, radius=0.1)
        add_rect(cur, LM, y, 0.14, hh, accent_dark)
        tf = add_box(cur, LM + 0.6, y + 0.25, CW - 1.2, hh - 0.45)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style_run(r, 26, True, hx(WHITE))
        est_h += hh + 0.1

    def add_code(text):
        nonlocal est_h
        ensure_body()
        code_lines = (text or "").split("\n")
        n_lines = len(code_lines)
        hh = min(n_lines * 0.42 + 0.35, 3.6)
        ensure_space(hh + 0.1)
        y = BODY_TOP + est_h
        add_rect(cur, LM, y, CW, hh, accent_dark, radius=0.08)
        tf = add_box(cur, LM + 0.4, y + 0.2, CW - 0.8, hh - 0.4)
        first = True
        for ln in code_lines[:8]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(4)
            r = p.add_run()
            r.text = ln if ln else " "
            style_run(r, 24, False, hx(WHITE), "Consolas")
        est_h += hh + 0.1

    def add_table_slide(title_text, headers, rows):
        """表格直接渲染在当前正文区（不再单独开页），避免标题页空白。"""
        nonlocal est_h
        from pptx.util import Inches as _In
        ensure_body()
        ncols = max(len(headers), max((len(rw) for rw in rows), default=0), 1)
        nrows = len(rows)
        vis_rows = min(nrows, 5)
        tbl_h = 0.7 + 0.62 * vis_rows
        ensure_space(tbl_h + 0.15)
        y = BODY_TOP + est_h
        tbl = cur.shapes.add_table(1 + vis_rows, ncols, _In(LM), _In(y), _In(CW), _In(tbl_h)).table
        colw = int(914400 * CW / ncols)
        for j in range(ncols):
            tbl.columns[j].width = colw
        for j in range(ncols):
            cell = tbl.cell(0, j)
            cell.text = headers[j] if j < len(headers) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent
            for pp in cell.text_frame.paragraphs:
                pp.alignment = PP_ALIGN.CENTER
                for rr in pp.runs:
                    style_run(rr, 24, True, hx(WHITE))
        for i, rw in enumerate(rows[:vis_rows]):
            for j in range(ncols):
                cell = tbl.cell(i + 1, j)
                cell.text = str(rw[j]) if j < len(rw) else ""
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = soft
                for pp in cell.text_frame.paragraphs:
                    pp.alignment = PP_ALIGN.CENTER
                    for rr in pp.runs:
                        style_run(rr, 24, False, hx(DARK_GRAY))
        if nrows > vis_rows:
            add_text(cur, LM, y + tbl_h + 0.12, CW, 0.5,
                     "… 共 %d 行数据" % nrows, size=24, color=MED_GRAY, align=PP_ALIGN.CENTER)
        est_h += tbl_h + 0.12

    def add_image(src, alt=""):
        nonlocal est_h
        ensure_body()
        local = resolve_image(src)
        if not local:
            para("图片不可用: " + str(src), size=24)
            return
        try:
            from PIL import Image
            im = Image.open(local)
            iw, ih = im.size
            max_w, max_h = 7.6, 2.9
            ratio = min(max_w / iw, max_h / ih) if iw and ih else 1
            w = max(1.0, iw * ratio)
            h = max(1.0, ih * ratio)
            need = h + 0.2
            ensure_space(need)
            y = BODY_TOP + est_h
            x = LM + (CW - w) / 2
            cur.shapes.add_picture(local, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            if alt:
                add_text(cur, LM, y + h + 0.05, CW, 0.4, alt, size=20, color=MED_GRAY, align=PP_ALIGN.CENTER)
                need += 0.5
            est_h += need
        except Exception as e:
            para("图片加载失败: %s" % e, size=24)

    def add_references(text):
        """参考文献条目：12pt 紧凑排版，逐条按每页容量自动切页，杜绝大段溢出。"""
        nonlocal est_h
        ensure_body()
        entries = [e for e in re.split(r"(?=\[\d+\])", text) if e.strip()]
        for e in entries:
            nl = est_lines(e, CW, 12)
            need = nl * 0.26 + 0.1 + 3 / 72.0
            ensure_space(need)
            p = body_tf.paragraphs[0] if not body_tf.paragraphs[0].runs and est_h == 0 else body_tf.add_paragraph()
            p.space_after = Pt(3)
            r = p.add_run()
            r.text = e
            style_run(r, 12, False, hx(DARK_GRAY))
            est_h += need

    sections = content_sections(spec)
    # 预扫描：每个 h1 后到下一个 h1 前的 h2 标题（供 PART 过渡页导航）；
    # 无 h2 的章节退回用首个要点块/段落作内容预告
    chapter_h2 = {}
    chapter_pv = {}
    cur_h1 = None
    for _s in sections:
        if _s.get("type") == "h1":
            cur_h1 = _s.get("text", "")
            chapter_h2.setdefault(cur_h1, [])
            chapter_pv.setdefault(cur_h1, [])
            continue
        if cur_h1 is None:
            continue
        if _s.get("type") == "h2":
            chapter_h2[cur_h1].append(_s.get("text", ""))
        elif not chapter_pv[cur_h1]:
            if _s.get("type") == "bullets":
                chapter_pv[cur_h1] = [it[:22] for it in _s.get("items", [])[:4]]
            elif _s.get("type") in ("p", "quote") and _s.get("text", "").strip():
                chapter_pv[cur_h1] = [_s["text"].strip()[:28]]
            elif _s.get("type") == "table":
                chapter_pv[cur_h1] = ["数据表：" + " · ".join((_s.get("headers") or [])[:3])]
            elif _s.get("type") == "image":
                chapter_pv[cur_h1] = [(_s.get("alt", "") or "章节配图")[:28]]
    for s in sections:
        typ = s.get("type")
        if typ == "h1":
            _h1t = s.get("text", "")
            nav = chapter_h2.get(_h1t) or chapter_pv.get(_h1t) or []
            open_slide(_h1t, is_divider=True, nav=nav)
        elif typ == "h2":
            if body_tf is None:
                open_slide(s.get("text", ""))
            elif est_h > 2.2:
                open_slide(s.get("text", ""))
            else:
                add_h2_row(s.get("text", ""))
        elif typ == "h3":
            add_h2_row(s.get("text", ""))
        elif typ == "p":
            if re.match(r"^\[\d+\]", s.get("text", "")):
                add_references(s.get("text", ""))
            else:
                para(s.get("text", ""), size=24)
        elif typ == "bullets":
            add_bullet_cards(s.get("items", []))
        elif typ == "numbered":
            add_numbered(s.get("items", []))
        elif typ == "quote":
            add_quote(s.get("text", ""))
        elif typ == "code":
            add_code(s.get("text", ""))
        elif typ == "table":
            add_table_slide(cur_title or "数据表", s.get("headers") or [], s.get("rows") or [])
        elif typ == "image":
            add_image(s.get("src", ""), s.get("alt", ""))
        elif typ == "divider":
            pass

    # ---------------- 页脚：页码 + 底部细条 ----------------
    total = len(slides) - 1
    for i, sl in enumerate(slides):
        if i == 0:
            continue
        add_hline(sl, 0, H - 0.08, W, accent, 2.5)
        add_text(sl, 12.0, 7.05, 1.1, 0.4, "%d/%d" % (i, total), size=16, color=MED_GRAY, align=PP_ALIGN.RIGHT)

    prs.save(out_path)
    _pptx_full_cleanup(out_path)


def _pptx_full_cleanup(outpath):
    """PPTX 后处理：移除所有 p:style 与主题阴影/3D，防止主题效果污染。
    模式借鉴 Mck-ppt-design-skill（Apache-2.0, github.com/likaku/Mck-ppt-design-skill）。"""
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree
    import zipfile, os
    tmppath = outpath + ".tmp"
    with zipfile.ZipFile(outpath, "r") as zin:
        with zipfile.ZipFile(tmppath, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith(".xml"):
                    root = _etree.fromstring(data)
                    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
                    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    for style in root.findall(".//{%s}style" % ns_p):
                        style.getparent().remove(style)
                    if "theme" in item.filename.lower():
                        for tag in ["outerShdw", "innerShdw", "scene3d", "sp3d"]:
                            for el in root.findall(".//{%s}%s" % (ns_a, tag)):
                                el.getparent().remove(el)
                    data = _etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)


def render_pptx(spec, out_path):
    """McKinsey 风格 PPT 渲染器（设计模式借鉴 Mck-ppt-design-skill，Apache-2.0，
    https://github.com/likaku/Mck-ppt-design-skill）：封面大标题 + 章节分隔页 +
    动作标题内容页 + 编号圆点小节 + 引用/代码/表格规范排版 + 页脚页码。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree

    t = theme(spec.get("theme"))

    def C(h):
        return RGBColor.from_string(str(h).lstrip("#"))

    accent = C(t["accent"])
    accent_dark = C(t["accent_dark"])
    light = C(t["light"])
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    MED_GRAY = RGBColor(0x66, 0x66, 0x66)
    LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    BG_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

    def hx(c):
        return "%02X%02X%02X" % (c[0], c[1], c[2])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = 13.333, 7.5
    LM = 0.8          # 原始英寸值，各绘图助手内部再 Inches() 换算
    CW = 11.733
    slides = []  # 所有页引用，结尾统一加页码

    def style_run(r, size=18, bold=False, color="333333", font="微软雅黑", italic=False):
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = RGBColor.from_string(str(color).lstrip("#"))
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)

    def clean(shape):
        sp = shape._element
        st = sp.find(qn("p:style"))
        if st is not None:
            sp.remove(st)

    def add_rect(slide, l, tp, w, h, color):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(tp), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        clean(sp)
        return sp

    def add_hline(slide, x, y, length, color, thickness=0.5):
        h = max(thickness * 12700, 6350)
        return add_rect(slide, x, y, length, h / 914400.0, color)

    def add_oval(slide, l, tp, size, letter, bg, fg=WHITE, font_size=15):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(tp), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = bg
        c.line.fill.background()
        c.shadow.inherit = False
        clean(c)
        tf = c.text_frame
        tf.paragraphs[0].text = letter
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.set("anchor", "ctr")
        for a in ["lIns", "tIns", "rIns", "bIns"]:
            bodyPr.set(a, "0")
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = fg
            r.font.name = "微软雅黑"
        return c

    def add_text(slide, l, tp, w, h, text, size=14, bold=False, color=DARK_GRAY,
                 font="微软雅黑", align=PP_ALIGN.LEFT, anchor="t"):
        tb = slide.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.set("anchor", anchor)
        for a in ["lIns", "tIns", "rIns", "bIns"]:
            bodyPr.set(a, "45720")
        lines = text if isinstance(text, list) else [text]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if i > 0:
                p.space_before = Pt(4)
            r = p.add_run()
            r.text = line
            style_run(r, size, bold, hx(color), font)
        return tb

    def add_box(slide, l, tp, w, h):
        tb = slide.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    def add_slide():
        s = prs.slides.add_slide(blank)
        slides.append(s)
        return s

    def add_page_number(slide, num, total):
        add_text(slide, 12.15, 7.08, 1.0, 0.3, "%d/%d" % (num, total),
                 size=9, color=MED_GRAY, align=PP_ALIGN.RIGHT)

    # ---- 封面页（大标题左对齐 + 顶部细线 + 底部粗线，McKinsey 风格） ----
    s = add_slide()
    add_rect(s, 0, 0, W, 0.06, accent)
    title = spec_title(spec) or "文档标题"
    add_text(s, 1.0, 1.85, 11.3, 1.4, title, size=40, bold=True, color=accent_dark)
    yy = 3.0
    if spec.get("subtitle"):
        add_text(s, 1.0, yy, 11.3, 0.6, spec["subtitle"], size=20, color=DARK_GRAY)
        yy += 0.85
    add_hline(s, 1.0, 6.62, 4.6, accent, 2.0)
    meta = []
    if spec.get("author"):
        meta.append(spec["author"])
    if spec.get("date"):
        meta.append(spec["date"])
    if meta:
        add_text(s, 1.0, 6.82, 11.3, 0.4, "　".join(meta), size=12, color=MED_GRAY)

    # ---- 内容结构 ----
    cur = None
    cur_title = ""
    body_tf = None
    est_h = 0.0
    BODY_TOP = 1.55
    BODY_MAX = 5.55
    h2_num = 0
    part_num = 0

    def open_slide(title_text, is_divider=False):
        nonlocal cur, cur_title, body_tf, est_h, part_num
        cur = add_slide()
        cur_title = title_text
        est_h = 0.0
        if is_divider:
            # 章节分隔页：左侧主题色竖条 + PART 标签 + 大标题
            part_num += 1
            add_rect(cur, 0, 0, 0.6, H, accent)
            add_text(cur, 1.25, 2.05, 10.0, 0.5, "PART %02d" % part_num, size=16, color=MED_GRAY)
            add_text(cur, 1.25, 2.65, 10.8, 1.2, title_text, size=28, bold=True, color=accent_dark)
            body_tf = None
        else:
            # 内容页：动作标题（底对齐）+ 细分隔线
            add_text(cur, LM, 0.32, 11.7, 0.6, title_text, size=22, bold=True, color=accent_dark, anchor="b")
            add_hline(cur, LM, 1.06, CW, LINE_GRAY, 0.75)
            body_tf = add_box(cur, LM, BODY_TOP, CW, BODY_MAX)
        return cur

    def ensure_body():
        nonlocal cur
        if body_tf is None:
            open_slide(cur_title)

    def ensure_space(need):
        nonlocal est_h
        # 只有确实放不下才续页（+0.02in 容差），避免短段落触发假性续页
        if est_h + need > BODY_MAX + 0.02:
            open_slide(cur_title + "（续）")

    def para(text="", size=16, bold=False, color="333333", before=0, after=6, bullet=None,
             num=None, italic=False, font="微软雅黑"):
        nonlocal est_h
        ensure_body()
        # 按真实折行行数计费：行高 ≈ 字号×1.3，另加段后距，最小一行的实际高度
        n_lines = _wrap_lines(text, CW, size)
        need = max(n_lines * (size / 72.0) * 1.3, size / 72.0) + after / 72.0 + 0.02
        ensure_space(need)
        p = body_tf.paragraphs[0] if not body_tf.paragraphs[0].runs and est_h == 0 else body_tf.add_paragraph()
        p.space_before = Pt(before)
        p.space_after = Pt(after)
        if bullet:
            rb = p.add_run()
            rb.text = "▪  "
            style_run(rb, size, True, hx(accent))
        if num is not None:
            rn = p.add_run()
            rn.text = "%d. " % num
            style_run(rn, size, True, hx(accent))
        if text:
            rt = p.add_run()
            rt.text = text
            style_run(rt, size, bold, color, font, italic)
        est_h += need

    def add_h2_row(text):
        nonlocal est_h, h2_num
        ensure_body()
        n_lines = _wrap_lines(text, CW - 0.62, 18)
        need = n_lines * 0.34 + 0.08
        ensure_space(need)
        h2_num += 1
        y = BODY_TOP + est_h
        add_oval(cur, LM, y - 0.03, 0.42, str(h2_num), accent)
        tf = add_box(cur, LM + 0.62, y - 0.07, CW - 0.62, n_lines * 0.34 + 0.1)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style_run(r, 18, True, hx(accent_dark))
        est_h += need

    def add_h3(text):
        nonlocal est_h
        ensure_body()
        n_lines = _wrap_lines(text, CW - 0.24, 16)
        need = n_lines * 0.30 + 0.08
        ensure_space(need)
        y = BODY_TOP + est_h
        add_rect(cur, LM, y + 0.06, 0.07, 0.26, accent)
        tf = add_box(cur, LM + 0.24, y - 0.04, CW - 0.24, n_lines * 0.30 + 0.08)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style_run(r, 16, True, "444444")
        est_h += need

    def add_quote(text):
        nonlocal est_h
        ensure_body()
        n_lines = _wrap_lines(text, CW - 0.6, 15)
        need = max(n_lines * 0.30 + 0.3, 1.0)
        ensure_space(need)
        y = BODY_TOP + est_h
        add_rect(cur, LM, y, CW, need - 0.15, light)
        add_rect(cur, LM, y, 0.07, need - 0.15, accent)
        tf = add_box(cur, LM + 0.3, y + 0.1, CW - 0.6, need - 0.35)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        style_run(r, 15, False, "555555", italic=True)
        est_h += need

    def add_code(text):
        nonlocal est_h
        ensure_body()
        lines = (text or "").split("\n")
        h = min(0.28 * len(lines) + 0.2, 3.2)
        ensure_space(h + 0.1)
        y = BODY_TOP + est_h
        add_rect(cur, LM, y, CW, h, BG_GRAY)
        tf = add_box(cur, LM + 0.25, y + 0.12, CW - 0.5, h - 0.2)
        first = True
        for ln in lines[:16]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = ln if ln else " "
            style_run(r, 12, False, "444444", "Consolas")
        est_h += h + 0.1

    def add_table_slide(title_text, headers, rows):
        s2 = add_slide()
        add_text(s2, LM, 0.32, 11.7, 0.6, title_text, size=22, bold=True, color=accent_dark, anchor="b")
        add_hline(s2, LM, 1.06, CW, LINE_GRAY, 0.75)
        ncols = max(len(headers), max((len(rw) for rw in rows), default=0), 1)
        nrows = len(rows)
        colw = CW / ncols
        hdr_y = 1.4
        for j in range(ncols):
            add_text(s2, LM + colw * j, hdr_y, colw, 0.42, headers[j] if j < len(headers) else "",
                     size=14, bold=True, color=accent_dark)
        add_hline(s2, LM, hdr_y + 0.46, CW, accent, 1.25)
        row_start = hdr_y + 0.6
        avail = 7.0 - row_start
        row_h = min(0.62, avail / nrows) if nrows else 0.62
        row_font = 13 if row_h >= 0.5 else 11
        for i, rw in enumerate(rows):
            ry = row_start + row_h * i
            for j in range(ncols):
                add_text(s2, LM + colw * j + 0.06, ry + 0.02, colw - 0.12, row_h - 0.04,
                         str(rw[j]) if j < len(rw) else "", size=row_font,
                         color=DARK_GRAY, anchor="ctr")
            add_hline(s2, LM, ry + row_h, CW, LINE_GRAY, 0.25)

    def add_image_slide(title_text, src, alt=""):
        s2 = add_slide()
        add_text(s2, LM, 0.32, 11.7, 0.6, title_text, size=22, bold=True, color=accent_dark, anchor="b")
        add_hline(s2, LM, 1.06, CW, LINE_GRAY, 0.75)
        local = resolve_image(src)
        if local:
            try:
                from PIL import Image
                im = Image.open(local)
                iw, ih = im.size
                max_w, max_h = 10.5, 5.2
                ratio = min(max_w / iw, max_h / ih) if iw and ih else 1
                w = max(1.0, iw * ratio)
                h = max(1.0, ih * ratio)
                left = LM + (CW - w) / 2
                top = 1.4 + (5.4 - h) / 2
                s2.shapes.add_picture(local, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
            except Exception as e:
                add_text(s2, LM, 3.0, CW, 0.8, "图片加载失败: %s" % e, size=14, color="C0392B")
        else:
            add_text(s2, LM, 3.0, CW, 0.8, "图片不可用: " + str(src), size=14, color="C0392B")
        if alt:
            add_text(s2, LM, 6.75, CW, 0.4, alt, size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)

    sections = content_sections(spec)
    for s in sections:
        typ = s.get("type")
        if typ == "h1":
            open_slide(s.get("text", ""), is_divider=True)
        elif typ == "h2":
            if body_tf is None:
                open_slide(s.get("text", ""))
            elif est_h > 2.2:
                open_slide(s.get("text", ""))
            else:
                add_h2_row(s.get("text", ""))
        elif typ == "h3":
            add_h3(s.get("text", ""))
        elif typ == "p":
            para(s.get("text", ""), size=16, before=2, after=8)
        elif typ == "bullets":
            for it in s.get("items", []):
                para(str(it), size=16, bullet=True, after=6)
        elif typ == "numbered":
            for n, it in enumerate(s.get("items", []), 1):
                para(str(it), size=16, num=n, after=6)
        elif typ == "quote":
            add_quote(s.get("text", ""))
        elif typ == "code":
            add_code(s.get("text", ""))
        elif typ == "table":
            add_table_slide(cur_title + "（附表）" if cur_title else "数据表", s.get("headers") or [], s.get("rows") or [])
        elif typ == "image":
            add_image_slide((cur_title + "（配图）") if cur_title else (s.get("alt") or "配图"), s.get("src", ""), s.get("alt", ""))
        elif typ == "divider":
            pass

    # ---- 页脚页码（封面除外） ----
    total = len(slides) - 1
    for i, sl in enumerate(slides):
        if i == 0:
            continue
        add_page_number(sl, i, total)

    prs.save(out_path)
    _pptx_full_cleanup(out_path)


def pptx_replace_text(shape, find, repl):
    changed = False
    if shape.shape_type == 6:
        for sub in shape.shapes:
            changed = pptx_replace_text(sub, find, repl) or changed
        return changed
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            full = "".join(r.text for r in p.runs)
            if find in full:
                new = full.replace(find, repl)
                if p.runs:
                    p.runs[0].text = new
                    for r in p.runs[1:]:
                        r.text = ""
                changed = True
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    full = "".join(r.text for r in p.runs)
                    if find in full:
                        new = full.replace(find, repl)
                        if p.runs:
                            p.runs[0].text = new
                            for r in p.runs[1:]:
                                r.text = ""
                        changed = True
    return changed


def edit_pptx(spec, in_path, out_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree

    t = theme(spec.get("theme"))
    prs = Presentation(in_path)
    ops = spec.get("ops") or []
    W, H = 13.333, 7.5

    def style_run(r, size=18, bold=False, color="333333", font="微软雅黑", italic=False):
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = RGBColor.from_string(str(color).lstrip("#"))
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)

    for op in ops:
        typ = op.get("type")
        if typ == "replace":
            find = op.get("find", "")
            repl = op.get("replace", "")
            if not find:
                continue
            for slide in prs.slides:
                for shape in slide.shapes:
                    pptx_replace_text(shape, find, repl)
        elif typ == "set-title":
            title = (op.get("title") or "").strip()
            if not title or not prs.slides:
                continue
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_text_frame and (shape.is_placeholder or "title" in (shape.name or "").lower()):
                    p = shape.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = title
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        r = p.add_run()
                        r.text = title
                    break
        elif typ == "append":
            # 追加为新幻灯片
            blank = prs.slide_layouts[6]
            for s in content_sections(op):
                if s.get("type") in ("h1", "h2"):
                    slide = prs.slides.add_slide(blank)
                    tf = slide.shapes.add_textbox(Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.85))
                    tf.word_wrap = True
                    p = tf.text_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = s.get("text", "")
                    style_run(r, 27, True, t["accent_dark"])
                elif s.get("type") == "p":
                    slide = prs.slides.add_slide(blank)
                    tf = slide.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(11.85), Inches(6.0))
                    tf.word_wrap = True
                    p = tf.text_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = s.get("text", "")
                    style_run(r, 18, False, "333333")
                elif s.get("type") == "bullets":
                    slide = prs.slides.add_slide(blank)
                    tf = slide.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(11.85), Inches(6.0))
                    tf.word_wrap = True
                    first = True
                    for it in s.get("items", []):
                        p = tf.text_frame.paragraphs[0] if first else tf.text_frame.add_paragraph()
                        first = False
                        rb = p.add_run()
                        rb.text = "▪  "
                        style_run(rb, 17, True, t["accent"])
                        rt = p.add_run()
                        rt.text = str(it)
                        style_run(rt, 17, False, "333333")
                elif s.get("type") == "image":
                    slide = prs.slides.add_slide(blank)
                    local = resolve_image(s.get("src", ""))
                    if local:
                        try:
                            from PIL import Image
                            im = Image.open(local)
                            iw, ih = im.size
                            max_w, max_h = 10.5, 5.4
                            ratio = min(max_w / iw, max_h / ih) if iw and ih else 1
                            w = max(1.0, iw * ratio)
                            h = max(1.0, ih * ratio)
                            left = (W - w) / 2
                            top = 1.0 + (5.5 - h) / 2
                            slide.shapes.add_picture(local, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
                        except Exception as e:
                            tf = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.85), Inches(0.8))
                            p = tf.text_frame.paragraphs[0]
                            r = p.add_run()
                            r.text = "图片加载失败: %s" % e
                            style_run(r, 18, False, "C0392B")
                    else:
                        tf = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.85), Inches(0.8))
                        p = tf.text_frame.paragraphs[0]
                        r = p.add_run()
                        r.text = "图片不可用: " + str(s.get("src", ""))
                        style_run(r, 18, False, "C0392B")
                    if s.get("alt"):
                        tf = slide.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(11.85), Inches(0.5))
                        p = tf.text_frame.paragraphs[0]
                        r = p.add_run()
                        r.text = str(s.get("alt", ""))
                        style_run(r, 16, False, "666666")
        elif typ == "restyle":
            accent = (op.get("accent") or "").strip()
            if not re.fullmatch(r"#[0-9A-Fa-f]{6}", accent):
                continue
            old_colors = set()
            for th in THEMES.values():
                old_colors.add(th["accent"].lstrip("#").upper())
                old_colors.add(th["accent_dark"].lstrip("#").upper())
            new_rgb = RGBColor.from_string(accent.lstrip("#"))
            for slide in prs.slides:
                for shape in slide.shapes:
                    try:
                        if shape.fill.type is not None and shape.fill.type == 1:  # MSO_FILL_TYPE.SOLID
                            fg = shape.fill.fore_color
                            rgb = fg.rgb
                            if rgb is not None and str(rgb).upper() in old_colors:
                                fg.rgb = new_rgb
                    except Exception:
                        pass
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                try:
                                    c = r.font.color
                                    rgb = c.rgb
                                    if rgb is not None and str(rgb).upper() in old_colors:
                                        c.rgb = new_rgb
                                except Exception:
                                    pass
    prs.save(out_path)


# ===========================================================================
# XLSX（Excel）
# ===========================================================================
def render_xlsx(spec, out_path):
    """从 sections 生成 Excel：h1/h2 作为分组标题行，table 写入工作表，bullets/paragraph 附注。"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    t = theme(spec.get("theme"))
    wb = Workbook()
    ws = wb.active
    ws.title = "数据"

    header_fill = PatternFill("solid", fgColor=t["accent"].lstrip("#"))
    header_font = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
    title_font = Font(name="微软雅黑", size=13, bold=True, color=t["accent_dark"].lstrip("#"))
    body_font = Font(name="微软雅黑", size=10.5)
    alt_fill = PatternFill("solid", fgColor=t["soft"].lstrip("#"))
    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    row_idx = 1
    if spec.get("title"):
        c = ws.cell(row=row_idx, column=1, value=spec["title"])
        c.font = Font(name="微软雅黑", size=16, bold=True, color=t["accent_dark"].lstrip("#"))
        ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=8)
        c.alignment = Alignment(horizontal="left", vertical="center")
        ws.row_dimensions[row_idx].height = 30
        row_idx += 2

    for s in content_sections(spec):
        typ = s.get("type")
        if typ == "table":
            headers = s.get("headers") or []
            rows = s.get("rows") or []
            if not headers and rows:
                headers = rows[0]
                rows = rows[1:]
            for j, h in enumerate(headers, 1):
                c = ws.cell(row=row_idx, column=j, value=h)
                c.fill = header_fill
                c.font = header_font
                c.alignment = Alignment(horizontal="center", vertical="center")
                c.border = border
            ws.row_dimensions[row_idx].height = 22
            row_idx += 1
            for i, r in enumerate(rows):
                for j in range(1, max(len(headers), len(r), 1) + 1):
                    val = r[j - 1] if j - 1 < len(r) else ""
                    c = ws.cell(row=row_idx, column=j, value=val)
                    c.font = body_font
                    c.border = border
                    if i % 2 == 1:
                        c.fill = alt_fill
                    c.alignment = Alignment(vertical="center", wrap_text=True)
                ws.row_dimensions[row_idx].height = 20
                row_idx += 1
            row_idx += 1
        elif typ in ("h1", "h2", "h3"):
            c = ws.cell(row=row_idx, column=1, value=s.get("text", ""))
            c.font = title_font
            c.alignment = Alignment(vertical="center")
            ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=8)
            ws.row_dimensions[row_idx].height = 24
            row_idx += 1
        elif typ in ("p", "quote"):
            c = ws.cell(row=row_idx, column=1, value=s.get("text", ""))
            c.font = body_font
            ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=8)
            c.alignment = Alignment(wrap_text=True, vertical="top")
            ws.row_dimensions[row_idx].height = max(18, 16 * (len(s.get("text", "")) // 60 + 1))
            row_idx += 1
        elif typ == "bullets":
            for it in s.get("items", []):
                c = ws.cell(row=row_idx, column=1, value="• " + str(it))
                c.font = body_font
                ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=8)
                row_idx += 1
        elif typ == "numbered":
            for n, it in enumerate(s.get("items", []), 1):
                c = ws.cell(row=row_idx, column=1, value="%d. %s" % (n, it))
                c.font = body_font
                ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=8)
                row_idx += 1
        elif typ == "divider":
            row_idx += 1

    # 自动列宽（估算）
    for col in range(1, 9):
        width = 12
        for r in range(1, min(ws.max_row, 200) + 1):
            v = ws.cell(row=r, column=col).value
            if v is not None:
                width = max(width, min(len(str(v)) * 1.9 + 4, 60))
        ws.column_dimensions[get_column_letter(col)].width = width

    wb.save(out_path)


def edit_xlsx(spec, in_path, out_path):
    """修改 Excel：replace=查找替换单元格文本；append-sheet=新增工作表（table 数据）；set-title=改第一个工作表名。"""
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    t = theme(spec.get("theme"))
    wb = load_workbook(in_path)
    ops = spec.get("ops") or []
    for op in ops:
        typ = op.get("type")
        if typ == "replace":
            find = op.get("find", "")
            repl = op.get("replace", "")
            if not find:
                continue
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for c in row:
                        if c.value is not None and find in str(c.value):
                            c.value = str(c.value).replace(find, repl)
        elif typ == "set-title":
            title = (op.get("title") or "").strip()
            if title and wb.worksheets:
                wb.worksheets[0].title = title
        elif typ == "append-sheet":
            ws = wb.create_sheet(title=(op.get("sheet") or "Sheet"))
            header_fill = PatternFill("solid", fgColor=t["accent"].lstrip("#"))
            header_font = Font(name="微软雅黑", size=11, bold=True, color="FFFFFF")
            body_font = Font(name="微软雅黑", size=10.5)
            thin = Side(style="thin", color="CCCCCC")
            border = Border(left=thin, right=thin, top=thin, bottom=thin)
            headers = op.get("headers") or []
            rows = op.get("rows") or []
            for j, h in enumerate(headers, 1):
                c = ws.cell(row=1, column=j, value=h)
                c.fill = header_fill
                c.font = header_font
                c.alignment = Alignment(horizontal="center", vertical="center")
                c.border = border
            for i, r in enumerate(rows, 2):
                for j in range(1, max(len(headers), len(r), 1) + 1):
                    val = r[j - 1] if j - 1 < len(r) else ""
                    c = ws.cell(row=i, column=j, value=val)
                    c.font = body_font
                    c.border = border
    wb.save(out_path)


# ===========================================================================
# PDF
# ===========================================================================
def render_pdf(spec, out_path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor, white
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (BaseDocTemplate, PageTemplate, Frame, Paragraph,
                                    Spacer, Table, TableStyle, HRFlowable, KeepTogether,
                                    PageBreak, NextPageTemplate)
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT

    t = theme(spec.get("theme"))
    FONT = "WQY"
    pdfmetrics.registerFont(TTFont(FONT, CJK_PDF_FONT, subfontIndex=0))
    try:
        pdfmetrics.registerFont(TTFont("MONO", MONO_PDF_FONT))
    except Exception:
        pass

    A = t["accent"]
    AD = t["accent_dark"]
    LIGHT = t["light"]
    SOFT = t["soft"]

    st_title = ParagraphStyle("title", fontName=FONT, fontSize=25, leading=32, alignment=TA_CENTER, textColor=HexColor(AD))
    st_sub = ParagraphStyle("sub", fontName=FONT, fontSize=13, leading=18, alignment=TA_CENTER, textColor=HexColor("#666666"))
    st_meta = ParagraphStyle("meta", fontName=FONT, fontSize=10.5, leading=15, alignment=TA_CENTER, textColor=HexColor("#888888"))
    st_h1 = ParagraphStyle("h1", fontName=FONT, fontSize=16.5, leading=22, textColor=HexColor(A), spaceBefore=14, spaceAfter=8)
    st_h2 = ParagraphStyle("h2", fontName=FONT, fontSize=13.5, leading=18, textColor=HexColor(AD), spaceBefore=10, spaceAfter=5)
    st_h3 = ParagraphStyle("h3", fontName=FONT, fontSize=11.5, leading=16, textColor=HexColor("#444444"), spaceBefore=8, spaceAfter=4)
    st_body = ParagraphStyle("body", fontName=FONT, fontSize=10.5, leading=17, alignment=TA_JUSTIFY, firstLineIndent=21, spaceAfter=5)
    st_bullet = ParagraphStyle("bullet", fontName=FONT, fontSize=10.5, leading=16.5, leftIndent=18, bulletIndent=6, spaceAfter=4)
    st_quote = ParagraphStyle("quote", fontName=FONT, fontSize=10.5, leading=16.5, textColor=HexColor("#666666"), leftIndent=22, spaceAfter=6)
    st_code = ParagraphStyle("code", fontName="MONO" if "MONO" in pdfmetrics.getRegisteredFontNames() else FONT,
                             fontSize=9, leading=13, leftIndent=16, rightIndent=16, backColor=HexColor(LIGHT), borderPadding=6, spaceAfter=4)

    title = spec_title(spec) or "文档"

    def on_page(canvas, doc_):
        canvas.saveState()
        canvas.setFont(FONT, 8.5)
        canvas.setFillColor(HexColor("#999999"))
        canvas.drawString(20 * mm, 13 * mm, title[:60])
        canvas.drawCentredString(A4[0] / 2.0, 13 * mm, "第 %d 页" % doc_.page)
        canvas.setStrokeColor(HexColor("#DDDDDD"))
        canvas.setLineWidth(0.5)
        canvas.line(20 * mm, 17 * mm, A4[0] - 20 * mm, 17 * mm)
        canvas.restoreState()

    def on_cover(canvas, doc_):
        canvas.saveState()
        canvas.setStrokeColor(HexColor(A))
        canvas.setLineWidth(2.5)
        canvas.line(20 * mm, A4[1] - 16 * mm, A4[0] - 20 * mm, A4[1] - 16 * mm)
        canvas.restoreState()

    doc = BaseDocTemplate(out_path, pagesize=A4,
                          leftMargin=20 * mm, rightMargin=20 * mm,
                          topMargin=24 * mm, bottomMargin=22 * mm)
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="main")
    doc.addPageTemplates([
        PageTemplate(id="cover", frames=[frame], onPage=on_cover),
        PageTemplate(id="body", frames=[frame], onPage=on_page),
    ])

    story = []

    def esc(text):
        return str(text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    # ---- 封面 ----
    story.append(Spacer(1, 90))
    if title:
        story.append(Paragraph(esc(title), st_title))
    story.append(Spacer(1, 14))
    story.append(HRFlowable(width="100%", thickness=1.2, color=HexColor(A), spaceBefore=2, spaceAfter=14))
    if spec.get("subtitle"):
        story.append(Paragraph(esc(spec["subtitle"]), st_sub))
        story.append(Spacer(1, 10))
    meta = []
    if spec.get("author"):
        meta.append(esc(spec["author"]))
    if spec.get("date"):
        meta.append(esc(spec["date"]))
    if meta:
        story.append(Paragraph("　".join(meta), st_meta))
    story.append(NextPageTemplate("body"))
    story.append(PageBreak())
    story.append(Spacer(1, 4))

    # ---- 正文 ----
    def render_section(s):
        typ = s.get("type")
        if typ == "h1":
            story.append(KeepTogether([Paragraph(esc(s.get("text", "")), st_h1),
                                       HRFlowable(width="100%", thickness=1.0, color=HexColor(A), spaceAfter=8)]))
        elif typ == "h2":
            story.append(KeepTogether([Paragraph(esc(s.get("text", "")), st_h2)]))
        elif typ == "h3":
            story.append(KeepTogether([Paragraph(esc(s.get("text", "")), st_h3)]))
        elif typ == "p":
            story.append(Paragraph(esc(s.get("text", "")), st_body))
        elif typ == "quote":
            q = Table([[Paragraph(esc(s.get("text", "")), st_quote)]], colWidths=[doc.width - 12])
            q.setStyle(TableStyle([
                ("LINEBEFORE", (0, 0), (0, -1), 2.5, HexColor(A)),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("BACKGROUND", (0, 0), (-1, -1), HexColor(SOFT)),
            ]))
            story.append(q)
            story.append(Spacer(1, 6))
        elif typ == "code":
            for line in (s.get("text", "") or "").split("\n"):
                story.append(Paragraph(esc(line if line else " "), st_code))
            story.append(Spacer(1, 6))
        elif typ == "bullets":
            for it in s.get("items", []):
                story.append(Paragraph(esc(it), st_bullet, bulletText="•"))
        elif typ == "numbered":
            for n, it in enumerate(s.get("items", []), 1):
                story.append(Paragraph(esc(it), st_bullet, bulletText="%d." % n))
        elif typ == "table":
            render_table(s)
        elif typ == "divider":
            story.append(HRFlowable(width="100%", thickness=0.6, color=HexColor("#CCCCCC"), spaceBefore=4, spaceAfter=10))

    def render_table(s):
        headers = [esc(h) for h in (s.get("headers") or [])]
        rows = [[esc(c) for c in r] for r in (s.get("rows") or [])]
        ncols = max(len(headers), max((len(r) for r in rows), default=0), 1)
        data = [headers] + rows
        tbl = Table(data, repeatRows=1)
        style = [
            ("BACKGROUND", (0, 0), (-1, 0), HexColor(A)),
            ("TEXTCOLOR", (0, 0), (-1, 0), white),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("FONTNAME", (0, 0), (-1, -1), FONT),
            ("FONTSIZE", (0, 0), (-1, -1), 9.5),
            ("LEADING", (0, 0), (-1, -1), 13),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ]
        for i in range(1, len(rows) + 1):
            if i % 2 == 0:
                style.append(("BACKGROUND", (0, i), (-1, i), HexColor(SOFT)))
        tbl.setStyle(TableStyle(style))
        story.append(KeepTogether([tbl]))
        story.append(Spacer(1, 8))

    for s in content_sections(spec):
        render_section(s)

    doc.build(story)


def edit_pdf(spec, in_path, out_path):
    info = extract_pdf(in_path)
    text = info["text"]
    ops = spec.get("ops") or []
    for op in ops:
        typ = op.get("type")
        if typ == "replace" and op.get("find"):
            text = text.replace(op["find"], op.get("replace", ""))
        elif typ == "set-title":
            pass
        elif typ == "append":
            extra = []
            for s in content_sections(op):
                if s.get("type") == "h1":
                    extra.append("# " + s.get("text", ""))
                elif s.get("type") == "h2":
                    extra.append("## " + s.get("text", ""))
                elif s.get("type") == "p":
                    extra.append(s.get("text", ""))
                elif s.get("type") == "bullets":
                    extra.extend("- " + str(it) for it in s.get("items", []))
                elif s.get("type") == "numbered":
                    extra.extend("%d. %s" % (i, it) for i, it in enumerate(s.get("items", []), 1))
            if extra:
                text = text + "\n\n" + "\n".join(extra)
    build_spec = {"title": spec.get("title") or "", "subtitle": spec.get("subtitle") or "",
                  "author": spec.get("author") or "", "date": spec.get("date") or "",
                  "theme": spec.get("theme") or "blue", "content": text}
    render_pdf(build_spec, out_path)


def edit_md(spec, in_path, out_path):
    with io.open(in_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    ops = spec.get("ops") or []
    for op in ops:
        typ = op.get("type")
        if typ == "replace" and op.get("find"):
            text = text.replace(op["find"], op.get("replace", ""))
        elif typ == "set-title" and op.get("title"):
            lines = text.split("\n")
            if lines and lines[0].startswith("#"):
                lines[0] = "# " + op["title"]
            else:
                lines.insert(0, "# " + op["title"])
            text = "\n".join(lines)
        elif typ == "append":
            extra = []
            for s in content_sections(op):
                if s.get("type") == "h1":
                    extra.append("# " + s.get("text", ""))
                elif s.get("type") == "h2":
                    extra.append("## " + s.get("text", ""))
                elif s.get("type") == "h3":
                    extra.append("### " + s.get("text", ""))
                elif s.get("type") == "p":
                    extra.append(s.get("text", ""))
                elif s.get("type") == "bullets":
                    extra.extend("- " + str(it) for it in s.get("items", []))
                elif s.get("type") == "numbered":
                    extra.extend("%d. %s" % (i, it) for i, it in enumerate(s.get("items", []), 1))
                elif s.get("type") == "quote":
                    extra.append("> " + s.get("text", ""))
                elif s.get("type") == "table":
                    headers = s.get("headers") or []
                    rows = s.get("rows") or []
                    extra.append("| " + " | ".join(headers) + " |")
                    extra.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    for r in rows:
                        extra.append("| " + " | ".join(str(c) for c in r) + " |")
            if extra:
                text = text.rstrip() + "\n\n" + "\n".join(extra) + "\n"
    with io.open(out_path, "w", encoding="utf-8") as f:
        f.write(text)


# ===========================================================================
# 文献检索与核查（PubMed/PMC E-utilities + Crossref，真实来源二次核查）
# ===========================================================================
def _urlencode(s):
    import urllib.parse
    return urllib.parse.quote_plus(str(s))


def _http_get(url, timeout=30):
    import urllib.request
    req = urllib.request.Request(url, headers={
        "User-Agent": "Mozilla/5.0 (docflow-literature/1.0; mailto:research@example.com)",
        "Accept": "application/json,text/plain;q=0.9,*/*;q=0.5",
    })
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read().decode("utf-8", errors="replace")


EUTILS = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"


def pubmed_esearch(term, retmax=8, mindate=None, maxdate=None, sort="relevance"):
    """PubMed 检索 → PMID 列表"""
    q = "term=%s&retmode=json&retmax=%d&sort=%s" % (_urlencode(term), retmax, _urlencode(sort))
    if mindate:
        q += "&mindate=%s" % _urlencode(mindate)
    if maxdate:
        q += "&maxdate=%s" % _urlencode(maxdate)
    url = "%s/esearch.fcgi?db=pubmed&%s" % (EUTILS, q)
    try:
        j = json.loads(_http_get(url))
        ids = (j.get("esearchresult") or {}).get("idlist") or []
        count = int((j.get("esearchresult") or {}).get("count", "0"))
        return ids, count
    except Exception as e:
        return [], -1


def pubmed_esummary(pmids):
    """PMID 列表 → 元数据摘要"""
    if not pmids:
        return {}
    url = "%s/esummary.fcgi?db=pubmed&retmode=json&id=%s" % (EUTILS, ",".join(pmids))
    try:
        j = json.loads(_http_get(url))
        return (j.get("result") or {})
    except Exception:
        return {}


def pubmed_efetch_abstract(pmids):
    """PMID 列表 → 纯文本摘要（含标题/作者/期刊行）"""
    if not pmids:
        return {}
    url = "%s/efetch.fcgi?db=pubmed&retmode=text&rettype=abstract&id=%s" % (EUTILS, ",".join(pmids))
    try:
        txt = _http_get(url)
    except Exception:
        return {}
    # 按 PMID 分块（efetch 文本以 1. / 2. 编号分隔）
    parts = {}
    blocks = re.split(r"\n(?=\d+\.\s)", txt)
    for b in blocks:
        m = re.match(r"^(\d+)\.\s", b)
        if m and int(m.group(1)) - 1 < len(pmids):
            parts[pmids[int(m.group(1)) - 1]] = b.strip()
    return parts


def _clean_abstract(text):
    if not text:
        return ""
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    # 去掉标题/作者/来源行（一般在前 3-4 行）
    out = []
    for l in lines:
        if re.match(r"^(Author information|PMID:|DOI:|PMCID:|Conflict of interest|©|Copyright)", l, re.I):
            continue
        out.append(l)
    return "\n".join(out).strip()


def _fmt_authors_et_al(names, max_authors=3):
    """GB/T 7714：作者 1, 作者 2, 作者 3, 等."""
    names = [n for n in names if n]
    if not names:
        return "佚名"
    if len(names) <= max_authors:
        return ", ".join(names)
    return ", ".join(names[:max_authors]) + ", 等"


def gb7714_citation(title, authors, journal, year, volume, issue, pages, doi=None, lang="en"):
    """生成 GB/T 7714-2015 顺序编码制引文（中华口腔医学会格式）"""
    au = _fmt_authors_et_al(authors)
    t = title or "无题"
    j = journal or "未知期刊"
    y = year or ""
    if lang == "zh":
        return "%s. %s[J]. %s, %s, %s(%s): %s." % (au, t, j, y, volume or "", issue or "", pages or "")
    # 英文：作者. 题名[J]. 刊名缩写, 年, 卷(期): 页码.
    base = "%s. %s[J]. %s, %s" % (au, t, j, y)
    vol_issue = ""
    if volume:
        vol_issue = volume
        if issue:
            vol_issue += "(%s)" % issue
    if vol_issue:
        base += ", %s" % vol_issue
    if pages:
        base += ": %s" % pages
    base += "."
    if doi:
        base += " DOI: %s." % doi
    return base


def lit_search(term, retmax=8):
    """检索 PubMed（真实来源），返回结构化条目列表"""
    ids, count = pubmed_esearch(term, retmax=retmax)
    if not ids:
        return {"count": count, "items": []}
    time.sleep(0.35)  # NCBI 限速 3 req/s
    result = pubmed_esummary(ids)
    time.sleep(0.35)
    abstracts = pubmed_efetch_abstract(ids)
    items = []
    for pid in ids:
        r = result.get(pid) or {}
        if not r:
            continue
        authors = []
        for a in (r.get("authors") or []):
            nm = a.get("name") or ""
            if nm:
                authors.append(nm.replace(",", " "))
        title = r.get("title") or ""
        journal = r.get("fulljournalname") or r.get("source") or ""
        year = ""
        pd = r.get("pubdate") or ""
        m = re.search(r"(\d{4})", pd)
        if m:
            year = m.group(1)
        doi = r.get("elocationid") or ""
        if doi.startswith("doi:"):
            doi = doi[4:]
        items.append({
            "pmid": pid,
            "title": title,
            "authors": authors,
            "journal": journal,
            "year": year,
            "volume": r.get("volume") or "",
            "issue": r.get("issue") or "",
            "pages": r.get("pages") or "",
            "doi": doi,
            "pubtype": r.get("pubtype") or [],
            "abstract": _clean_abstract(abstracts.get(pid, "")),
            "citation": gb7714_citation(title, authors, journal, year,
                                        r.get("volume") or "", r.get("issue") or "",
                                        r.get("pages") or "", doi or None),
            "source": "pubmed",
        })
    return {"count": count, "items": items}


def crossref_search(query, rows=5):
    """Crossref 检索（DOI 权威）"""
    url = "https://api.crossref.org/works?query.bibliographic=%s&rows=%d&select=DOI,title,author,container-title,issued,volume,issue,page,type,URL" % (_urlencode(query), rows)
    try:
        j = json.loads(_http_get(url))
    except Exception as e:
        return {"error": str(e), "items": []}
    items = []
    # 只保留正式文献类型，过滤 decision letter / 评论等噪声
    GOOD_TYPES = {"journal-article", "proceedings-article", "book-chapter", "book", "dissertation"}
    BAD_URL_MARK = "/decision"
    for w in (j.get("message") or {}).get("items") or []:
        wtype = (w.get("type") or "").lower()
        url = w.get("URL") or ""
        if wtype and wtype not in GOOD_TYPES:
            continue
        if BAD_URL_MARK in url:
            continue
        title = (w.get("title") or [""])[0]
        journal = (w.get("container-title") or [""])[0]
        issued = w.get("issued") or {}
        year = ""
        dp = (issued.get("date-parts") or [[None]])[0]
        if dp and dp[0]:
            year = str(dp[0])
        authors = []
        for a in (w.get("author") or []):
            nm = " ".join([a.get("family") or "", a.get("given") or ""]).strip()
            if nm:
                authors.append(nm)
        items.append({
            "doi": w.get("DOI") or "",
            "title": title,
            "authors": authors,
            "journal": journal,
            "year": year,
            "volume": w.get("volume") or "",
            "issue": w.get("issue") or "",
            "pages": w.get("page") or "",
            "type": wtype,
            "citation": gb7714_citation(title, authors, journal, year,
                                        w.get("volume") or "", w.get("issue") or "",
                                        w.get("page") or "", w.get("DOI") or None),
            "source": "crossref",
        })
    return {"items": items}


def _norm(s):
    return re.sub(r"[^a-z0-9\u4e00-\u9fff]", "", (s or "").lower())


def lit_verify(references):
    """对用户提供的引文逐条核查真实性：PubMed/Crossref 双重比对"""
    results = []
    for ref in references:
        title = (ref.get("title") or "").strip()
        doi = (ref.get("doi") or "").strip()
        pmid = (ref.get("pmid") or "").strip()
        authors = ref.get("authors") or []
        year = ref.get("year") or ""
        entry = {"input": ref, "verified": False, "matched": None, "notes": []}

        # 1) DOI 精确核查
        if doi:
            try:
                j = json.loads(_http_get("https://api.crossref.org/works/%s" % _urlencode(doi)))
                w = j.get("message") or {}
                t = (w.get("title") or [""])[0]
                entry["matched"] = {
                    "doi": doi,
                    "title": t,
                    "authors": [(" ".join([a.get("family") or "", a.get("given") or ""])).strip() for a in (w.get("author") or [])],
                    "journal": (w.get("container-title") or [""])[0],
                    "year": str(((w.get("issued") or {}).get("date-parts") or [[None]])[0][0] or ""),
                    "volume": w.get("volume") or "",
                    "issue": w.get("issue") or "",
                    "pages": w.get("page") or "",
                }
                entry["matched"]["citation"] = gb7714_citation(t, entry["matched"]["authors"], entry["matched"]["journal"],
                                                               entry["matched"]["year"], entry["matched"]["volume"],
                                                               entry["matched"]["issue"], entry["matched"]["pages"], doi)
                entry["verified"] = True
                entry["notes"].append("DOI 在 Crossref 真实存在")
            except Exception:
                entry["notes"].append("DOI 未能在 Crossref 解析")

        # 2) PMID 核查
        if pmid and not entry["verified"]:
            result = pubmed_esummary([pmid])
            r = result.get(pmid)
            if r:
                entry["verified"] = True
                entry["matched"] = {
                    "pmid": pmid,
                    "title": r.get("title") or "",
                    "authors": [a.get("name", "").replace(",", " ") for a in (r.get("authors") or [])],
                    "journal": r.get("fulljournalname") or r.get("source") or "",
                    "year": re.search(r"(\d{4})", r.get("pubdate") or "").group(1) if re.search(r"(\d{4})", r.get("pubdate") or "") else "",
                    "volume": r.get("volume") or "",
                    "issue": r.get("issue") or "",
                    "pages": r.get("pages") or "",
                }
                entry["matched"]["citation"] = gb7714_citation(entry["matched"]["title"], entry["matched"]["authors"],
                                                              entry["matched"]["journal"], entry["matched"]["year"],
                                                              entry["matched"]["volume"], entry["matched"]["issue"],
                                                              entry["matched"]["pages"])
                entry["notes"].append("PMID 在 PubMed 真实存在")
            else:
                entry["notes"].append("PMID 未能在 PubMed 找到")

        # 3) 标题比对（PubMed 精确检索）
        if title and not entry["verified"]:
            time.sleep(0.35)
            ids, _cnt = pubmed_esearch('"%s"[Title]' % title, retmax=5)
            if ids:
                time.sleep(0.35)
                result = pubmed_esummary(ids)
                for pid in ids:
                    r = result.get(pid) or {}
                    if not r:
                        continue
                    rt = r.get("title") or ""
                    if _norm(rt) == _norm(title) or (_norm(rt) and _norm(rt)[:30] == _norm(title)[:30]):
                        entry["verified"] = True
                        entry["matched"] = {
                            "pmid": pid,
                            "title": rt,
                            "authors": [a.get("name", "").replace(",", " ") for a in (r.get("authors") or [])],
                            "journal": r.get("fulljournalname") or r.get("source") or "",
                            "year": re.search(r"(\d{4})", r.get("pubdate") or "").group(1) if re.search(r"(\d{4})", r.get("pubdate") or "") else "",
                            "volume": r.get("volume") or "",
                            "issue": r.get("issue") or "",
                            "pages": r.get("pages") or "",
                        }
                        entry["matched"]["citation"] = gb7714_citation(rt, entry["matched"]["authors"], entry["matched"]["journal"],
                                                                      entry["matched"]["year"], entry["matched"]["volume"],
                                                                      entry["matched"]["issue"], entry["matched"]["pages"])
                        entry["notes"].append("标题在 PubMed 精确匹配 PMID %s" % pid)
                        break
            if not entry["verified"]:
                entry["notes"].append("标题未在 PubMed 精确匹配（可能为中文库或不存在）")

        # 4) Crossref 兜底
        if title and not entry["verified"]:
            cr = crossref_search(title, rows=3)
            for c in cr.get("items") or []:
                if _norm(c["title"]) == _norm(title) or (_norm(c["title"]) and _norm(c["title"])[:30] == _norm(title)[:30]):
                    entry["verified"] = True
                    entry["matched"] = c
                    entry["notes"].append("标题在 Crossref 精确匹配")
                    break
            if not entry["verified"]:
                entry["notes"].append("Crossref 未找到精确匹配")

        results.append(entry)
    return {"results": results}


# ---------------------------------------------------------------------------
# 中药/方剂数据核验（tcm-verify）
# ---------------------------------------------------------------------------
# 内置权威核验规则：剂量上限（参照《中国药典》2020 常用剂量范围）、
# 十八反十九畏配伍禁忌、毒性药警示。
# 等级：A=《中国药典》原文级 / B=教材共识级 / C=经典医籍级（未在线核验）
DOSE_LIMITS = {
    "甘草": 10, "黄连": 5, "黄芩": 10, "栀子": 10, "知母": 12, "黄柏": 12,
    "细辛": 3, "附子": 15, "半夏": 9, "吴茱萸": 5, "山豆根": 6, "青黛": 3,
    "人参": 9, "西洋参": 6, "五味子": 6, "乌梅": 12, "酸枣仁": 15, "远志": 10,
    "石菖蒲": 10, "木通": 6, "滑石": 20, "泽泻": 10, "车前子": 15, "苍术": 9,
    "厚朴": 10, "枳实": 10, "香附": 10, "郁金": 10, "莪术": 9, "三棱": 10,
    "桃仁": 10, "红花": 10, "三七": 9, "川芎": 10, "丹参": 15, "益母草": 30,
    "白花蛇舌草": 60, "半枝莲": 30, "紫花地丁": 30, "鱼腥草": 25, "蒲公英": 15,
    "板蓝根": 15, "金银花": 15, "连翘": 15, "生石膏": 60, "茯苓": 15, "白术": 12,
    "黄芪": 30, "党参": 30, "太子参": 30, "山药": 30, "薏苡仁": 30, "砂仁": 6,
    "陈皮": 10, "肉桂": 5, "干姜": 10, "吴茱萸": 5, "小茴香": 6, "肉豆蔻": 10,
    "诃子": 10, "赤石脂": 12, "玄明粉": 9, "硼砂": 3, "冰片": 0.3, "儿茶": 3,
    "血竭": 2, "乳香": 5, "没药": 5, "土茯苓": 60, "忍冬藤": 30, "路路通": 10,
    "独活": 10, "防己": 10, "木香": 6, "槟榔": 10, "乌药": 10, "荔枝核": 10,
    "山楂": 12, "神曲": 15, "麦芽": 15, "鸡内金": 10, "莱菔子": 12, "地榆": 15,
    "槐花": 10, "仙茅": 10, "淫羊藿": 10, "巴戟天": 10, "肉苁蓉": 10, "菟丝子": 12,
    "补骨脂": 10, "杜仲": 10, "续断": 15, "骨碎补": 9, "益智仁": 10, "牛膝": 12,
    "王不留行": 10, "苍耳子": 10, "辛夷": 10, "荆芥": 10, "防风": 10, "羌活": 10,
    "紫苏叶": 10, "生姜": 10, "桔梗": 10, "射干": 9, "桑白皮": 12, "葶苈子": 10,
    "紫菀": 10, "款冬花": 10, "百部": 9, "川贝母": 10, "浙贝母": 10, "瓜蒌": 15,
    "白蔹": 10, "红藤": 15, "败酱草": 15, "白头翁": 15, "秦皮": 12, "马齿苋": 15,
    "芦根": 30, "天花粉": 15, "桑叶": 10, "菊花": 10, "薄荷": 6, "牛蒡子": 12,
    "升麻": 10, "葛根": 15, "白芷": 10, "柴胡": 10, "当归": 12, "白芍": 15,
    "赤芍": 12, "生地黄": 15, "熟地黄": 15, "麦冬": 12, "玄参": 15, "牡丹皮": 12,
    "夏枯草": 15, "决明子": 15, "女贞子": 12, "墨旱莲": 12, "山茱萸": 12,
    "枸杞子": 12, "肉桂": 5, "石斛": 12, "玉竹": 12, "黄精": 15, "白鲜皮": 10,
    "地肤子": 15, "苦参": 9, "紫草": 10, "白及": 15, "仙鹤草": 12, "白茅根": 30,
    "侧柏叶": 12, "炮姜": 6, "龙胆草": 6, "马勃": 6, "山慈菇": 9, "威灵仙": 10,
    "桑寄生": 15, "僵蚕": 10, "天麻": 10, "钩藤": 12, "龙骨": 30, "牡蛎": 30,
    "珍珠母": 25, "酸枣仁": 15, "远志": 10, "合欢皮": 12, "夜交藤": 15,
    "茯苓": 15, "泽泻": 10, "茵陈": 15, "藿香": 10, "佩兰": 10, "苍术": 9,
    "木香": 6, "川芎": 10, "丹参": 15, "桃仁": 10, "红花": 10, "三七": 9,
    "延胡索": 10, "川楝子": 10, "枸杞": 12, "杜仲": 10, "续断": 15, "菟丝子": 12,
}

TOXIC_HERBS = {
    "附子": "有毒；先煎久煎（30-60分钟）；孕妇禁用",
    "半夏": "生品有毒，内服宜制；不宜与乌头类同用",
    "细辛": "有小毒；用量不宜过大（≤3g）",
    "苍耳子": "有毒；过量可致中毒",
    "山豆根": "有毒；用量不宜过大",
    "吴茱萸": "有小毒",
    "仙茅": "有毒",
    "槟榔": "多服久服可致中毒",
    "苦参": "用量过大可致中毒",
    "青黛": "难溶于水，一般入丸散",
    "木通": "关木通有肾毒性（现药典已删）；须用川木通/木通，肾功不全者禁用",
    "雄黄": "有毒；不入煎剂",
    "朱砂": "有毒；不入煎剂；不宜久服",
    "马钱子": "大毒；炮制后入丸散",
    "斑蝥": "大毒",
    "巴豆": "大毒；制霜入丸散",
    "甘遂": "有毒；不宜与甘草同用",
    "大戟": "有毒；不宜与甘草同用",
    "芫花": "有毒；不宜与甘草同用",
    "牵牛子": "有毒；不宜与巴豆同用",
    "藜芦": "有毒；不宜与人参、沙参、丹参、玄参、苦参、细辛、芍药同用",
}

SHIBA_FAN = ["甘草反海藻", "甘草反京大戟", "甘草反红大戟", "甘草反甘遂", "甘草反芫花",
             "乌头反半夏", "乌头反瓜蒌", "乌头反贝母", "乌头反白蔹", "乌头反白及",
             "藜芦反人参", "藜芦反沙参", "藜芦反丹参", "藜芦反玄参", "藜芦反苦参",
             "藜芦反细辛", "藜芦反芍药"]
SHIJIU_WEI = ["硫黄畏朴硝", "水银畏砒霜", "狼毒畏密陀僧", "巴豆畏牵牛", "丁香畏郁金",
              "川乌畏犀角", "草乌畏犀角", "牙硝畏三棱", "官桂畏赤石脂", "人参畏五灵脂"]


def tcm_query_data():
    """加载 medkit 中医药数据库并返回元信息（供查询工具缓存校验）。"""
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "medkit", "tcm_data.json")
    if not os.path.exists(path):
        err("tcm_data.json 不存在: %s" % path)
    with io.open(path, "r", encoding="utf-8") as f:
        d = json.load(f)
    return {
        "herbs": len(d.get("herbs") or []),
        "pairs": len(d.get("pairs") or []),
        "formulas": len(d.get("formulas") or []),
        "incompatibilities": (d.get("incompatibilities") or {}).get("十八反") or [],
        "path": path,
    }


def tcm_verify(items, with_lit=True):
    """核验中药/方剂数据条目。
    items: [{kind: 'herb'|'pair'|'formula'|'incompatibility', name, dose?, composition?, notes?}]
    输出：每条核验结论（内置规则 + 可选 PubMed 佐证检索）。
    """
    results = []
    for it in items:
        kind = it.get("kind") or "herb"
        name = (it.get("name") or "").strip()
        entry = {"kind": kind, "name": name, "checks": [], "risk": "ok", "verified": False}

        if kind == "herb":
            dose = it.get("dose") or ""
            # 剂量上限核验
            m = re.findall(r"(\d+(?:\.\d+)?)\s*~\s*(\d+(?:\.\d+)?)", dose) or re.findall(r"(\d+(?:\.\d+)?)-(\d+(?:\.\d+)?)", dose)
            limit = DOSE_LIMITS.get(name)
            if m and limit is not None:
                hi = float(m[0][1])
                if hi > limit:
                    entry["checks"].append("⚠ 剂量上限 %s g 超过参考上限 %s g（《中国药典》2020 常用量）" % (hi, limit))
                    entry["risk"] = "dose-over"
                else:
                    entry["checks"].append("✅ 剂量上限 %s g 在参考范围（≤%s g）内" % (hi, limit))
                    entry["verified"] = True
            elif not m:
                entry["checks"].append("⚠ 未能解析剂量范围，需人工核对")
                entry["risk"] = "dose-unparsed"
            # 毒性药核验
            tox = TOXIC_HERBS.get(name)
            if tox:
                entry["checks"].append("☣ 毒性药：%s" % tox)
                entry["risk"] = entry["risk"] if entry["risk"] != "ok" else "toxic"
        elif kind == "pair":
            comp = (it.get("composition") or name or "").replace("、", "-").replace("，", "-").replace(" ", "")
            # 十八反十九畏核验：检查组合内是否出现反药对
            fan_hit = None
            for f in SHIBA_FAN + SHIJIU_WEI:
                pair_names = f.split("反") if "反" in f else f.split("畏")
                if len(pair_names) == 2:
                    a, b = pair_names[0].strip(), pair_names[1].strip()
                    if a in comp and b in comp:
                        fan_hit = f
                        break
            if fan_hit:
                entry["checks"].append("❌ 配伍禁忌：%s" % fan_hit)
                entry["risk"] = "incompatible"
            else:
                entry["checks"].append("✅ 未检出十八反/十九畏禁忌组合")
                entry["verified"] = True
        elif kind == "formula":
            comp = (it.get("composition") or "").replace("、", "-").replace("，", "-").replace(" ", "")
            fan_hit = None
            for f in SHIBA_FAN + SHIJIU_WEI:
                pair_names = f.split("反") if "反" in f else f.split("畏")
                if len(pair_names) == 2:
                    a, b = pair_names[0].strip(), pair_names[1].strip()
                    if a in comp and b in comp:
                        fan_hit = f
                        break
            if fan_hit:
                entry["checks"].append("❌ 方剂内配伍禁忌：%s" % fan_hit)
                entry["risk"] = "incompatible"
            else:
                entry["checks"].append("✅ 未检出方内十八反/十九畏冲突")
                entry["verified"] = True
        elif kind == "incompatibility":
            entry["checks"].append("ℹ 十八反/十九畏为《中国药典》2020 用药禁忌（等级 A）")
            entry["verified"] = True

        # 文献佐证（PubMed，可选；每个条目最多 1 次检索，控制速率）
        if with_lit and kind in ("herb", "pair") and entry["risk"] != "incompatible":
            try:
                time.sleep(0.35)
                term = name.replace("-", " ") + " traditional Chinese medicine"
                ids, _cnt = pubmed_esearch(term, retmax=3)
                if ids:
                    entry["literature"] = ids[:3]
                    entry["checks"].append("📚 PubMed 检索到 %d 条相关文献（PMID: %s）" % (len(ids), ", ".join(ids[:3])))
            except Exception:
                pass

        results.append(entry)
    return {"results": results}






    data = sys.stdin.buffer.read()
    data = re.sub(rb"\s+", b"", data)
    try:
        raw = base64.b64decode(data, validate=True)
    except Exception:
        err("base64 解码失败")
    with open(out_path, "wb") as f:
        f.write(raw)
    ok({"size": len(raw)})


def cmd_decode_file(b64_path, out_path):
    """从 base64 文本文件解码写二进制（绕开 stdin 传递的不确定性）。"""
    if not os.path.exists(b64_path):
        err("base64 文件不存在: %s" % b64_path)
    with io.open(b64_path, "r", encoding="utf-8", errors="replace") as f:
        data = f.read()
    data = re.sub(r"\s+", "", data)
    try:
        raw = base64.b64decode(data, validate=True)
    except Exception as e:
        err("base64 解码失败: %s" % e)
    with open(out_path, "wb") as f:
        f.write(raw)
    ok({"size": len(raw)})


def cmd_extract(path):
    if not os.path.exists(path):
        err("文件不存在: %s" % path)
    info = extract_any(path)
    text = info.get("text", "") or ""
    info["chars"] = len(text)
    info["size"] = os.path.getsize(path)
    info.setdefault("pages", 0)
    info.setdefault("slides", 0)
    ok(info)


def cmd_meta(path):
    if not os.path.exists(path):
        err("文件不存在: %s" % path)
    fmt = ext_of(path)
    ok({"format": fmt, "size": os.path.getsize(path), "name": os.path.basename(path)})


def cmd_image_search(query, max_results=5, image_type="all"):
    res = image_search(query, max_results=max_results, image_type=image_type)
    if res.get("error"):
        err(res["error"])
    ok(res)


def cmd_image_download(url, out_path):
    if not url:
        err("缺少图片 URL")
    local = resolve_image(url)
    if not local:
        err("图片下载失败: %s" % url)
    # resolve_image 可能已下载到 tmp，若指定输出则复制过去
    if os.path.abspath(local) != os.path.abspath(out_path):
        import shutil
        shutil.copyfile(local, out_path)
    ok({"path": out_path, "size": os.path.getsize(out_path)})


def cmd_image_info(path):
    if not os.path.exists(path):
        err("图片不存在: %s" % path)
    try:
        ok(image_info(path))
    except Exception as e:
        err("无法读取图片: %s" % e)


def cmd_image_recognize(path):
    if not os.path.exists(path):
        err("图片不存在: %s" % path)
    ok(image_recognize(path))


def _load_spec(argv_spec_path=None):
    """优先从 JSON 文件读 spec，否则读 stdin（向后兼容）。"""
    if argv_spec_path:
        if not os.path.exists(argv_spec_path):
            err("spec 文件不存在: %s" % argv_spec_path)
        with io.open(argv_spec_path, "r", encoding="utf-8") as f:
            return json.loads(f.read() or "{}")
    return json.loads(sys.stdin.read() or "{}")


def cmd_create(fmt, out_path, spec_path=None):
    spec = _load_spec(spec_path)
    fmt = fmt.lower()
    if fmt in ("ppt",):
        err("不支持旧版二进制 .ppt 格式，请先在 Office 中另存为 .pptx")
    if fmt == "docx":
        render_docx(spec, out_path)
    elif fmt == "pptx":
        if spec.get("visual") or spec.get("style") == "visual":
            render_pptx_visual(spec, out_path)
        else:
            render_pptx(spec, out_path)
    elif fmt == "pdf":
        render_pdf(spec, out_path)
    elif fmt == "md":
        text = spec.get("content") or ""
        with io.open(out_path, "w", encoding="utf-8") as f:
            f.write(text)
    elif fmt == "txt":
        sections = content_sections(spec)
        lines = []
        for s in sections:
            typ = s.get("type")
            if typ in ("h1", "h2", "h3"):
                lines.append("#" * int(typ[1]) + " " + s.get("text", ""))
            elif typ == "p":
                lines.append(s.get("text", ""))
            elif typ == "bullets":
                lines.extend("- " + str(it) for it in s.get("items", []))
            elif typ == "numbered":
                lines.extend("%d. %s" % (i, it) for i, it in enumerate(s.get("items", []), 1))
            elif typ == "quote":
                lines.append("> " + s.get("text", ""))
            elif typ == "table":
                headers = s.get("headers") or []
                lines.append(" | ".join(headers))
                for r in s.get("rows", []):
                    lines.append(" | ".join(str(c) for c in r))
            lines.append("")
        with io.open(out_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
    elif fmt == "xlsx":
        render_xlsx(spec, out_path)
    else:
        err("不支持的输出格式: %s" % fmt)
    if not os.path.exists(out_path):
        err("生成失败：未产生输出文件")
    ok({"size": os.path.getsize(out_path)})


def cmd_edit(fmt, in_path, out_path, spec_path=None):
    spec = _load_spec(spec_path)
    fmt = fmt.lower()
    if fmt == "ppt":
        err("不支持旧版二进制 .ppt 格式，请先在 Office 中另存为 .pptx")
    if not os.path.exists(in_path):
        err("源文件不存在: %s" % in_path)
    if fmt == "docx":
        edit_docx(spec, in_path, out_path)
    elif fmt == "pptx":
        edit_pptx(spec, in_path, out_path)
    elif fmt == "pdf":
        edit_pdf(spec, in_path, out_path)
    elif fmt in ("md", "txt"):
        edit_md(spec, in_path, out_path)
    elif fmt == "xlsx":
        edit_xlsx(spec, in_path, out_path)
    else:
        err("不支持的编辑格式: %s" % fmt)
    ok({"size": os.path.getsize(out_path)})


def main():
    try:
        for stream in (sys.stdin, sys.stdout, sys.stderr):
            try:
                stream.reconfigure(encoding="utf-8")
            except Exception:
                pass
        # ── 多用户路径栅栏 ──
        # 插件对非管理员用户以 `python engine.py __as__ <用户根> <命令> ...` 调用；
        # 栅栏生效后，所有命令的文件路径必须位于用户根之内，否则直接拒绝
        # （防止分用户文档工作流经引擎读取主用户项目内容——后台引擎进程内兜底）。
        global FENCE_ROOT
        FENCE_ROOT = None
        if len(sys.argv) > 2 and sys.argv[1] == "__as__":
            FENCE_ROOT = os.path.realpath(sys.argv[2])
            sys.argv = [sys.argv[0]] + list(sys.argv[3:])
        if len(sys.argv) < 2:
            err("用法: docflow_engine.py decode|extract|create|edit|meta …")
        cmd = sys.argv[1]
        if cmd == "decode":
            if len(sys.argv) != 3:
                err("decode 需要 <out_path>")
            cmd_decode(_fence(sys.argv[2]))
        elif cmd == "decode-file":
            if len(sys.argv) != 4:
                err("decode-file 需要 <b64_path> <out_path>")
            cmd_decode_file(_fence(sys.argv[2]), _fence(sys.argv[3]))
        elif cmd == "extract":
            if len(sys.argv) != 3:
                err("extract 需要 <file>")
            cmd_extract(_fence(sys.argv[2]))
        elif cmd == "meta":
            if len(sys.argv) != 3:
                err("meta 需要 <file>")
            cmd_meta(_fence(sys.argv[2]))
        elif cmd == "create":
            if len(sys.argv) not in (4, 5):
                err("create 需要 <fmt> <out_path> [spec_path]")
            cmd_create(sys.argv[2], _fence(sys.argv[3]), _fence(sys.argv[4]) if len(sys.argv) == 5 else None)
        elif cmd == "edit":
            if len(sys.argv) not in (5, 6):
                err("edit 需要 <fmt> <in> <out> [spec_path]")
            cmd_edit(sys.argv[2], _fence(sys.argv[3]), _fence(sys.argv[4]), _fence(sys.argv[5]) if len(sys.argv) == 6 else None)
        elif cmd == "lit-search":
            if len(sys.argv) not in (3, 4):
                err("lit-search 需要 <term> [retmax]")
            retmax = int(sys.argv[3]) if len(sys.argv) == 4 else 8
            ok(lit_search(sys.argv[2], retmax=retmax))
        elif cmd == "lit-crossref":
            if len(sys.argv) not in (3, 4):
                err("lit-crossref 需要 <query> [rows]")
            rows = int(sys.argv[3]) if len(sys.argv) == 4 else 5
            ok(crossref_search(sys.argv[2], rows=rows))
        elif cmd == "lit-verify":
            if len(sys.argv) != 3:
                err("lit-verify 需要 <refs_json_path>")
            with io.open(_fence(sys.argv[2]), "r", encoding="utf-8") as f:
                refs = json.loads(f.read() or "[]")
            ok(lit_verify(refs))
        elif cmd == "image-search":
            if len(sys.argv) not in (3, 4, 5):
                err("image-search 需要 <query> [max_results] [photo|vector|all]")
            n = int(sys.argv[3]) if len(sys.argv) >= 4 else 5
            itype = sys.argv[4] if len(sys.argv) == 5 else "all"
            cmd_image_search(sys.argv[2], max_results=n, image_type=itype)
        elif cmd == "image-download":
            if len(sys.argv) != 4:
                err("image-download 需要 <url> <out_path>")
            cmd_image_download(sys.argv[2], _fence(sys.argv[3]))
        elif cmd == "image-info":
            if len(sys.argv) != 3:
                err("image-info 需要 <file>")
            cmd_image_info(_fence(sys.argv[2]))
        elif cmd == "image-recognize":
            if len(sys.argv) != 3:
                err("image-recognize 需要 <file>")
            cmd_image_recognize(_fence(sys.argv[2]))
        elif cmd == "tcm-verify":
            if len(sys.argv) not in (3, 4):
                err("tcm-verify 需要 <items_json_path> [with_lit]")
            with io.open(_fence(sys.argv[2]), "r", encoding="utf-8") as f:
                items = json.loads(f.read() or "[]")
            with_lit = (len(sys.argv) == 4 and sys.argv[3] == "0") is False
            ok(tcm_verify(items, with_lit=with_lit))
        elif cmd == "tcm-query":
            ok(tcm_query_data())
        elif cmd == "dir-mtimes":
            # 返回目录内全部文件的真实修改时间（毫秒 epoch），供插件扫描时兜底显示
            # （DSH fs.stat 不提供 mtime，shell stat 又受沙箱限制不可靠）
            if len(sys.argv) != 3:
                err("dir-mtimes 需要 <dir>")
            d = _fence(sys.argv[2])
            out = {}
            try:
                for name in os.listdir(d):
                    p = os.path.join(d, name)
                    if os.path.isfile(p):
                        out[p] = int(os.stat(p).st_mtime * 1000)
            except Exception:
                pass
            print(json.dumps(out, ensure_ascii=False))
        else:
            err("未知命令: %s" % cmd)
    except SystemExit:
        raise
    except Exception as e:
        err("内部错误: %s\n%s" % (e, traceback.format_exc()))


if __name__ == "__main__":
    import traceback
    main()
